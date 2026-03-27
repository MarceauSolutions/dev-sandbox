#!/usr/bin/env python3
"""
Template Manager for Interview Prep PowerPoint Generator

Allows users to:
1. Load an existing PowerPoint as a template
2. Extract content/structure from an existing presentation
3. Apply a new theme to an existing presentation
4. Create a session from an existing file for editing

Usage:
    # Load existing PPTX as template and create editing session
    python template_manager.py --load .tmp/my_presentation.pptx

    # Apply new theme to existing presentation
    python template_manager.py --load .tmp/my_presentation.pptx --apply-theme professional

    # Extract research data from existing presentation (for regeneration)
    python template_manager.py --load .tmp/my_presentation.pptx --extract-data

    # Copy template to new working file
    python template_manager.py --load .tmp/template.pptx --copy-to .tmp/new_presentation.pptx
"""

import argparse
import json
import sys
import os
import shutil
import re
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# Try to import session manager
try:
    from session_manager import create_session, get_active_session
    SESSION_ENABLED = True
except ImportError:
    SESSION_ENABLED = False


# Theme color schemes (same as pptx_generator.py)
THEMES = {
    "modern": {
        "primary": RGBColor(0x1a, 0x1a, 0x2e),
        "secondary": RGBColor(0x16, 0x21, 0x3e),
        "accent": RGBColor(0xe9, 0x4d, 0x60),
        "text": RGBColor(0xff, 0xff, 0xff),
        "text_dark": RGBColor(0x1a, 0x1a, 0x2e),
        "background": RGBColor(0xff, 0xff, 0xff),
    },
    "professional": {
        "primary": RGBColor(0x00, 0x33, 0x66),
        "secondary": RGBColor(0x00, 0x52, 0x99),
        "accent": RGBColor(0xff, 0x99, 0x00),
        "text": RGBColor(0xff, 0xff, 0xff),
        "text_dark": RGBColor(0x33, 0x33, 0x33),
        "background": RGBColor(0xf5, 0xf5, 0xf5),
    },
    "minimal": {
        "primary": RGBColor(0x2c, 0x3e, 0x50),
        "secondary": RGBColor(0x34, 0x49, 0x5e),
        "accent": RGBColor(0x27, 0xae, 0x60),
        "text": RGBColor(0xff, 0xff, 0xff),
        "text_dark": RGBColor(0x2c, 0x3e, 0x50),
        "background": RGBColor(0xff, 0xff, 0xff),
    }
}


def load_template(pptx_path: str) -> dict:
    """
    Load an existing PowerPoint and extract its structure.
    Returns metadata about the presentation.
    """
    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    try:
        prs = Presentation(str(path))

        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "title": "",
                "text_content": [],
                "has_images": False,
                "image_count": 0
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # First text is usually title
                        if not slide_data["title"]:
                            slide_data["title"] = text
                        else:
                            slide_data["text_content"].append(text[:200])

                # Check for images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    slide_data["has_images"] = True
                    slide_data["image_count"] += 1

            slides_info.append(slide_data)

        return {
            "success": True,
            "file_path": str(path.absolute()),
            "file_name": path.name,
            "slide_count": len(prs.slides),
            "slides": slides_info,
            "width": prs.slide_width,
            "height": prs.slide_height
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def extract_company_role(pptx_path: str) -> dict:
    """
    Try to extract company name and role from an existing presentation.
    Looks at title slide and common patterns.
    """
    template_info = load_template(pptx_path)
    if not template_info.get("success"):
        return template_info

    company = ""
    role = ""

    slides = template_info.get("slides", [])

    # Check first slide (usually title)
    if slides:
        title_slide = slides[0]
        title = title_slide.get("title", "")

        # Company is usually the main title
        if title:
            company = title

        # Role is usually in subtitle or second text element
        text_content = title_slide.get("text_content", [])
        for text in text_content:
            if "interview" in text.lower() or "preparation" in text.lower():
                # Extract role from "Interview Preparation: Role Name"
                match = re.search(r'(?:preparation|prep)[:\s]+(.+)', text, re.IGNORECASE)
                if match:
                    role = match.group(1).strip()
                    break
            elif ":" in text:
                # Could be "Role: Position Name"
                role = text.split(":")[-1].strip()
                break

    # Also check slide 6 which is usually "Role: X"
    for slide in slides:
        title = slide.get("title", "")
        if title.lower().startswith("role:"):
            role = title.replace("Role:", "").replace("role:", "").strip()
            break

    return {
        "success": True,
        "company": company,
        "role": role,
        "slide_count": template_info.get("slide_count"),
        "file_path": template_info.get("file_path")
    }


def copy_template(source_path: str, dest_path: str) -> dict:
    """
    Copy a template to a new working file.
    """
    source = Path(source_path)
    dest = Path(dest_path)

    if not source.exists():
        return {"success": False, "error": f"Source file not found: {source_path}"}

    try:
        shutil.copy(source, dest)
        return {
            "success": True,
            "source": str(source),
            "destination": str(dest),
            "message": f"Copied template to {dest.name}"
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def create_session_from_template(pptx_path: str, company: str = None, role: str = None) -> dict:
    """
    Create an editing session from an existing template.
    This allows the user to edit the template using pptx_editor.py
    """
    if not SESSION_ENABLED:
        return {"success": False, "error": "Session manager not available"}

    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    # Try to extract company/role if not provided
    if not company or not role:
        extracted = extract_company_role(pptx_path)
        if extracted.get("success"):
            company = company or extracted.get("company", "Unknown Company")
            role = role or extracted.get("role", "Unknown Role")

    # Create session
    session = create_session(
        company=company or "Template",
        role=role or "Editing Session",
        research_file=None,
        pptx_file=path.name,
        theme="custom"
    )

    return {
        "success": True,
        "session_id": session.get("session_id"),
        "pptx_file": path.name,
        "company": company,
        "role": role,
        "message": f"Session created for editing {path.name}"
    }


def list_available_templates(tmp_dir: str = ".tmp") -> dict:
    """
    List all PowerPoint files available as templates.
    """
    tmp_path = Path(tmp_dir)
    if not tmp_path.exists():
        return {"success": False, "error": f"Directory not found: {tmp_dir}"}

    templates = []
    for f in tmp_path.iterdir():
        if f.suffix.lower() in [".pptx", ".ppt"]:
            info = {
                "name": f.name,
                "path": str(f),
                "size": f.stat().st_size,
                "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat()
            }

            # Try to extract company/role
            extracted = extract_company_role(str(f))
            if extracted.get("success"):
                info["company"] = extracted.get("company", "")
                info["role"] = extracted.get("role", "")
                info["slide_count"] = extracted.get("slide_count", 0)

            templates.append(info)

    # Sort by modification time (newest first)
    templates.sort(key=lambda x: x.get("modified", ""), reverse=True)

    return {
        "success": True,
        "templates": templates,
        "count": len(templates)
    }


def main():
    parser = argparse.ArgumentParser(
        description="Template Manager for Interview Prep PowerPoint",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # List available templates
  python template_manager.py --list

  # Load a template and create editing session
  python template_manager.py --load .tmp/my_presentation.pptx

  # Copy template to new file
  python template_manager.py --load .tmp/template.pptx --copy-to .tmp/new_version.pptx

  # Extract company/role from existing presentation
  python template_manager.py --load .tmp/presentation.pptx --extract-info
        """
    )

    parser.add_argument("--list", "-l", action="store_true",
                        help="List available templates in .tmp/")
    parser.add_argument("--load", help="Path to PowerPoint template to load")
    parser.add_argument("--copy-to", help="Copy template to new file path")
    parser.add_argument("--extract-info", action="store_true",
                        help="Extract company/role info from template")
    parser.add_argument("--create-session", action="store_true",
                        help="Create editing session from template")
    parser.add_argument("--company", help="Company name (for session)")
    parser.add_argument("--role", help="Role name (for session)")
    parser.add_argument("--json", action="store_true",
                        help="Output in JSON format")

    args = parser.parse_args()

    # List templates
    if args.list:
        result = list_available_templates()
        if args.json:
            print(json.dumps(result, indent=2))
        else:
            if result.get("success"):
                print(f"\n{'=' * 60}")
                print(f"Available Templates ({result.get('count')} found)")
                print(f"{'=' * 60}\n")

                for t in result.get("templates", []):
                    company = t.get("company", "")
                    role = t.get("role", "")
                    slides = t.get("slide_count", "?")
                    print(f"  {t['name']}")
                    if company or role:
                        print(f"    └─ {company} | {role} | {slides} slides")
                    print(f"    └─ Modified: {t.get('modified', 'Unknown')[:19]}")
                    print()
            else:
                print(f"Error: {result.get('error')}")
        return 0

    # Load template
    if args.load:
        # Copy to new file
        if args.copy_to:
            result = copy_template(args.load, args.copy_to)
            if args.json:
                print(json.dumps(result, indent=2))
            else:
                if result.get("success"):
                    print(f"✅ {result.get('message')}")
                else:
                    print(f"❌ Error: {result.get('error')}")
            return 0 if result.get("success") else 1

        # Extract info
        if args.extract_info:
            result = extract_company_role(args.load)
            if args.json:
                print(json.dumps(result, indent=2))
            else:
                if result.get("success"):
                    print(f"\n📄 Template: {args.load}")
                    print(f"   Company: {result.get('company', 'Unknown')}")
                    print(f"   Role: {result.get('role', 'Unknown')}")
                    print(f"   Slides: {result.get('slide_count', 'Unknown')}")
                else:
                    print(f"❌ Error: {result.get('error')}")
            return 0 if result.get("success") else 1

        # Create session
        if args.create_session:
            result = create_session_from_template(
                args.load,
                company=args.company,
                role=args.role
            )
            if args.json:
                print(json.dumps(result, indent=2))
            else:
                if result.get("success"):
                    print(f"\n✅ Session created!")
                    print(f"   Session ID: {result.get('session_id')}")
                    print(f"   File: {result.get('pptx_file')}")
                    print(f"   Company: {result.get('company')}")
                    print(f"   Role: {result.get('role')}")
                    print(f"\n→ You can now use pptx_editor.py to edit this presentation")
                else:
                    print(f"❌ Error: {result.get('error')}")
            return 0 if result.get("success") else 1

        # Default: load and show info
        result = load_template(args.load)
        if args.json:
            print(json.dumps(result, indent=2))
        else:
            if result.get("success"):
                print(f"\n{'=' * 60}")
                print(f"Template Loaded: {result.get('file_name')}")
                print(f"{'=' * 60}")
                print(f"📊 Slides: {result.get('slide_count')}")
                print(f"\nSlide Structure:")

                for slide in result.get("slides", []):
                    num = slide.get("slide_number")
                    title = slide.get("title", "Untitled")[:50]
                    images = slide.get("image_count", 0)
                    img_str = f" [{images} image(s)]" if images > 0 else ""
                    print(f"  {num:2}. {title}{img_str}")

                print(f"\n→ To create editing session: python template_manager.py --load {args.load} --create-session")
            else:
                print(f"❌ Error: {result.get('error')}")
        return 0 if result.get("success") else 1

    parser.print_help()
    return 0


if __name__ == "__main__":
    sys.exit(main())
