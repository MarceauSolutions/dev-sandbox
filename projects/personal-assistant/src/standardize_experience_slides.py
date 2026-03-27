#!/usr/bin/env python3
"""
Standardize Experience Slides to Match Target Format

Takes a PowerPoint and reformats experience slides to match a consistent layout:
- Title at top (0.55", 0.3")
- Image on left (0.75", 1.5") - 5" x 5"
- Description on right (6.55", 1.5")
- Relevance section on right (6.55", 4.5")
- No accent bar
- Consistent colors matching the professional theme

Usage:
    python standardize_experience_slides.py --input .tmp/file.pptx --slides 14,15,16,17,18
"""

import argparse
import sys
import json
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# Professional theme colors (matching slide 22)
THEME = {
    "background": RGBColor(0x00, 0x33, 0x66),  # Navy blue
    "title_color": RGBColor(0xFF, 0xFF, 0xFF),  # White
    "text_color": RGBColor(0xFF, 0xFF, 0xFF),   # White
    "accent": RGBColor(0xFF, 0x99, 0x00),       # Orange
}

# Layout positions (matching slide 22)
LAYOUT = {
    "title": {"left": 0.55, "top": 0.3, "width": 12.23, "height": 0.54},
    "image": {"left": 0.75, "top": 1.5, "width": 5.0, "height": 5.0},
    "description": {"left": 6.55, "top": 1.5, "width": 5.9, "height": 2.5},
    "relevance": {"left": 6.55, "top": 4.5, "width": 5.9, "height": 2.0},
}


def extract_slide_content(slide):
    """Extract text content and image from a slide."""
    content = {
        "title": "",
        "description": "",
        "relevance": "",
        "image_blob": None,
        "image_ext": None
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()

            # Identify content by keywords or position
            if not content["title"] and shape.top < Inches(1):
                content["title"] = text
            elif "relevance" in text.lower() or "why it matters" in text.lower():
                # Extract the relevance text
                if ":" in text:
                    content["relevance"] = text.split(":", 1)[1].strip()
                else:
                    content["relevance"] = text
            elif text and not content["description"]:
                content["description"] = text

        # Extract image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                content["image_blob"] = shape.image.blob
                content["image_ext"] = shape.image.ext
            except:
                pass

    return content


def clear_slide_content(slide):
    """Remove all shapes except the background."""
    shapes_to_remove = []
    for shape in slide.shapes:
        # Keep the main background shape
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.width > Inches(13) and shape.height > Inches(7):
                continue
        shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=None):
    """Add a text box with specified formatting."""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color

    return textbox


def rebuild_experience_slide(slide, content, slide_num):
    """Rebuild a slide with the standardized format."""

    # Clear existing content (except background)
    shapes_to_remove = []
    background_shape = None

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.width > Inches(13) and shape.height > Inches(7):
                background_shape = shape
                continue
        shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Set background color if we have the background shape
    if background_shape:
        try:
            background_shape.fill.solid()
            background_shape.fill.fore_color.rgb = THEME["background"]
        except:
            pass

    # Add title
    title = content.get("title", f"Experience {slide_num}")
    # Clean up generic titles
    if "Relevant Experience" in title:
        title = content.get("description", "")[:50] if content.get("description") else f"Experience {slide_num}"
        if title and len(title) > 40:
            title = title[:40] + "..."

    title_box = add_text_box(
        slide, title,
        LAYOUT["title"]["left"], LAYOUT["title"]["top"],
        LAYOUT["title"]["width"], LAYOUT["title"]["height"],
        font_size=28, bold=True, color=THEME["title_color"]
    )

    # Add description
    desc = content.get("description", "")
    if desc:
        desc_box = add_text_box(
            slide, desc,
            LAYOUT["description"]["left"], LAYOUT["description"]["top"],
            LAYOUT["description"]["width"], LAYOUT["description"]["height"],
            font_size=14, color=THEME["text_color"]
        )

    # Add relevance section
    relevance = content.get("relevance", "")
    if relevance:
        relevance_text = f"Relevance:\n{relevance}"
        rel_box = add_text_box(
            slide, relevance_text,
            LAYOUT["relevance"]["left"], LAYOUT["relevance"]["top"],
            LAYOUT["relevance"]["width"], LAYOUT["relevance"]["height"],
            font_size=12, color=THEME["text_color"]
        )

    # Add image if we have one
    if content.get("image_blob"):
        import io
        image_stream = io.BytesIO(content["image_blob"])
        try:
            slide.shapes.add_picture(
                image_stream,
                Inches(LAYOUT["image"]["left"]),
                Inches(LAYOUT["image"]["top"]),
                width=Inches(LAYOUT["image"]["width"])
            )
        except Exception as e:
            print(f"  Warning: Could not add image to slide {slide_num}: {e}")

    return True


def standardize_slides(pptx_path: str, slide_numbers: list, output_path: str = None):
    """Standardize specified slides to match the target format."""

    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    try:
        prs = Presentation(str(path))
        total_slides = len(prs.slides)

        results = []

        for slide_num in slide_numbers:
            if slide_num < 1 or slide_num > total_slides:
                results.append({"slide": slide_num, "status": "skipped", "reason": "out of range"})
                continue

            slide = prs.slides[slide_num - 1]

            # Extract current content
            content = extract_slide_content(slide)

            # Rebuild with standardized format
            success = rebuild_experience_slide(slide, content, slide_num)

            results.append({
                "slide": slide_num,
                "status": "updated" if success else "failed",
                "title": content.get("title", "")[:40]
            })

        # Save
        output = output_path or str(path)
        prs.save(output)

        return {
            "success": True,
            "output": output,
            "slides_processed": len(slide_numbers),
            "results": results
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def main():
    parser = argparse.ArgumentParser(
        description="Standardize experience slides to match target format"
    )
    parser.add_argument("--input", required=True, help="Input PowerPoint file")
    parser.add_argument("--slides", required=True, help="Comma-separated slide numbers to standardize")
    parser.add_argument("--output", help="Output file (defaults to overwriting input)")
    parser.add_argument("--json", action="store_true", help="Output in JSON format")

    args = parser.parse_args()

    # Parse slide numbers
    try:
        slide_numbers = [int(x.strip()) for x in args.slides.split(",")]
    except ValueError:
        print("Error: --slides must be comma-separated numbers (e.g., 14,15,16)")
        return 1

    result = standardize_slides(args.input, slide_numbers, args.output)

    if args.json:
        print(json.dumps(result, indent=2))
    else:
        if result.get("success"):
            print(f"✅ Standardized {result['slides_processed']} slides")
            print(f"   Output: {result['output']}")
            for r in result.get("results", []):
                status = "✓" if r["status"] == "updated" else "✗"
                print(f"   {status} Slide {r['slide']}: {r.get('title', '')[:30]}")
        else:
            print(f"❌ Error: {result.get('error')}")
            return 1

    return 0


if __name__ == "__main__":
    sys.exit(main())
