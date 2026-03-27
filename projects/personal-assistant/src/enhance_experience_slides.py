#!/usr/bin/env python3
"""
Enhance Experience Slides Content

Updates experience slides (14-18) to match the professional format of slides 19-22
with detailed descriptions and specific relevance to the target role.

Usage:
    python enhance_experience_slides.py --input .tmp/file.pptx
    python enhance_experience_slides.py --input .tmp/file.pptx --slides 14,15,16
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# Navy theme colors
NAVY_BG = RGBColor(0x00, 0x33, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Enhanced content for Brookhaven slides 14-18
# Each entry: (title, description, relevance)
ENHANCED_CONTENT = {
    14: {
        "title": "Aerospace Oxygen Systems (Aerox)",
        "description": (
            "As Director of Engineering at Aerox, designed and certified aerospace oxygen "
            "systems for FAA/EASA regulated aircraft. Managed complete product lifecycle from "
            "concept through certification, including oxygen regulators, masks, and distribution "
            "systems operating at altitudes up to 45,000 feet with precise pressure control and "
            "fail-safe mechanisms."
        ),
        "relevance": (
            "Relevance:\n"
            "Oxygen system design requires the same precision engineering and safety-critical "
            "mindset needed for BNL scientific instruments. Experience with pressure regulation, "
            "flow control, and certification processes directly applies to designing reliable "
            "components for particle detectors and cryogenic systems."
        )
    },
    15: {
        "title": "AW609 VTOL Qualification Testing",
        "description": (
            "Led technical qualification testing for the Leonardo AW609 tiltrotor aircraft "
            "per DO-160 environmental testing standards. Coordinated vibration, thermal cycling, "
            "altitude, and EMI/EMC testing campaigns. Developed test fixtures and procedures "
            "for validating avionics and mechanical systems under extreme operating conditions."
        ),
        "relevance": (
            "Relevance:\n"
            "DO-160 qualification testing methodology parallels the rigorous validation required "
            "for scientific instruments at BNL. Experience designing test fixtures and analyzing "
            "environmental stress data translates directly to qualifying detector components for "
            "radiation, thermal, and vibration environments."
        )
    },
    16: {
        "title": "NASA 747-SP Inerting System",
        "description": (
            "Collaborated with NASA on the installation and flight testing of a prototype fuel "
            "tank inerting system aboard the 747-SP research aircraft. Designed mounting brackets, "
            "routing for nitrogen distribution lines, and integration with existing aircraft systems. "
            "Participated in flight test campaigns to validate system performance."
        ),
        "relevance": (
            "Relevance:\n"
            "Direct experience working with a government research organization on experimental "
            "apparatus installation. Understanding of how to integrate complex systems into "
            "existing infrastructure mirrors the challenge of installing new instruments at "
            "established BNL facilities like NSLS-II and RHIC."
        )
    },
    17: {
        "title": "Engineering Documentation Review",
        "description": (
            "Over 15 years of experience as engineering drawing checker and documentation reviewer "
            "across aerospace, defense, and industrial sectors. Ensured compliance with ASME Y14 "
            "standards, GD&T accuracy, and manufacturing feasibility. Mentored junior engineers "
            "on proper documentation practices and tolerance stack-up analysis."
        ),
        "relevance": (
            "Relevance:\n"
            "Scientific instruments require meticulous documentation for fabrication, assembly, "
            "and long-term maintenance. Experience reviewing complex drawings ensures designs are "
            "manufacturable and maintainable—critical for one-of-a-kind detector components that "
            "must perform reliably for decades."
        )
    },
    18: {
        "title": "GD&T and CAD Expertise",
        "description": (
            "Expert-level application of Geometric Dimensioning and Tolerancing per ASME Y14.5 "
            "standards. Proficient in SolidWorks, Creo/Pro-E, and CATIA for complex surface "
            "modeling, assembly design, and simulation. Created parametric models enabling rapid "
            "design iteration and manufacturing optimization."
        ),
        "relevance": (
            "Relevance:\n"
            "Precision scientific instruments demand tight tolerances and sophisticated 3D modeling. "
            "GD&T expertise ensures components fit and function correctly in complex assemblies. "
            "CAD proficiency enables rapid prototyping and iteration—essential for developing "
            "custom detector components and experimental apparatus at BNL."
        )
    }
}


def update_text_box(shape, new_text):
    """Update text box content while preserving formatting."""
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    # Clear existing paragraphs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    # Set new text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = new_text
        run.font.color.rgb = WHITE
        run.font.size = Pt(14)

    return True


def enhance_slide(slide, slide_num, content):
    """Enhance a single slide with new content."""
    changes = {"slide": slide_num, "updated": []}

    # Find text boxes by position
    description_box = None
    relevance_box = None
    title_box = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        left = shape.left / 914400  # Convert to inches
        top = shape.top / 914400

        # Title is near top (0.3")
        if top < 0.5:
            title_box = shape
        # Description is at (6.55", 1.5")
        elif 6.0 < left < 7.0 and 1.0 < top < 2.0:
            description_box = shape
        # Relevance is at (6.55", 4.5")
        elif 6.0 < left < 7.0 and 4.0 < top < 5.0:
            relevance_box = shape

    # Update description
    if description_box:
        tf = description_box.text_frame
        tf.clear()
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = content["description"]
        run.font.color.rgb = WHITE
        run.font.size = Pt(14)
        changes["updated"].append("description")

    # Update relevance
    if relevance_box:
        tf = relevance_box.text_frame
        tf.clear()
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = content["relevance"]
        run.font.color.rgb = WHITE
        run.font.size = Pt(14)
        changes["updated"].append("relevance")

    return changes


def enhance_experience_slides(pptx_path: str, slide_numbers: list = None, output_path: str = None):
    """
    Enhance experience slides with detailed content.

    Args:
        pptx_path: Path to input PPTX file
        slide_numbers: List of slide numbers to enhance (default: 14-18)
        output_path: Output file path (defaults to input)
    """
    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    try:
        prs = Presentation(str(path))
        total_slides = len(prs.slides)

        # Default to slides 14-18
        if slide_numbers is None:
            slide_numbers = [14, 15, 16, 17, 18]

        results = []

        for slide_num in slide_numbers:
            if slide_num not in ENHANCED_CONTENT:
                results.append({
                    "slide": slide_num,
                    "status": "skipped",
                    "reason": "no enhanced content defined"
                })
                continue

            if slide_num < 1 or slide_num > total_slides:
                results.append({
                    "slide": slide_num,
                    "status": "skipped",
                    "reason": "out of range"
                })
                continue

            slide = prs.slides[slide_num - 1]
            content = ENHANCED_CONTENT[slide_num]
            changes = enhance_slide(slide, slide_num, content)
            changes["status"] = "enhanced"
            results.append(changes)

        # Save
        output = output_path or str(path)
        prs.save(output)

        return {
            "success": True,
            "output": output,
            "total_slides": total_slides,
            "slides_enhanced": len([r for r in results if r.get("status") == "enhanced"]),
            "results": results
        }

    except Exception as e:
        import traceback
        return {"success": False, "error": str(e), "traceback": traceback.format_exc()}


def main():
    parser = argparse.ArgumentParser(
        description="Enhance experience slides with detailed professional content"
    )
    parser.add_argument("--input", required=True, help="Input PowerPoint file")
    parser.add_argument("--slides", help="Comma-separated slide numbers (default: 14-18)")
    parser.add_argument("--output", help="Output file (defaults to input)")
    parser.add_argument("--open", action="store_true", help="Open file after processing")

    args = parser.parse_args()

    # Parse slide numbers if provided
    slide_numbers = None
    if args.slides:
        try:
            slide_numbers = [int(x.strip()) for x in args.slides.split(",")]
        except ValueError:
            print("Error: --slides must be comma-separated numbers")
            return 1

    print(f"\n{'='*60}")
    print(f"Enhancing Experience Slides")
    print(f"{'='*60}\n")

    result = enhance_experience_slides(args.input, slide_numbers, args.output)

    if result.get("success"):
        print(f"✅ Enhanced {result['slides_enhanced']} slides")
        print(f"   Output: {result['output']}")

        print(f"\n📋 Changes by slide:")
        for r in result.get("results", []):
            if r.get("status") == "enhanced":
                updated = ", ".join(r.get("updated", []))
                print(f"   ✓ Slide {r['slide']}: {updated}")
            else:
                print(f"   ✗ Slide {r['slide']}: {r.get('reason', 'skipped')}")
    else:
        print(f"❌ Error: {result.get('error')}")
        if result.get("traceback"):
            print(result["traceback"])
        return 1

    # Open if requested
    if args.open:
        import os
        os.system(f'open "{result.get("output", args.input)}"')

    return 0


if __name__ == "__main__":
    sys.exit(main())
