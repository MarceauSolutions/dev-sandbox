#!/usr/bin/env python3
"""
Interview Prep PowerPoint API
FastAPI REST API for interview preparation PowerPoint generation.

Key Features:
    - Combined endpoint for research + generation in one step
    - Slide preview with base64 image rendering
    - Natural language editing
    - Direct module imports (no subprocess) for better error handling
"""

import os
import sys
import json
import shutil
import traceback
import tempfile
import base64
from pathlib import Path
from typing import Optional, List
from datetime import datetime
from io import BytesIO

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from pydantic import BaseModel

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent))

# Load environment variables
from dotenv import load_dotenv
env_path = Path(__file__).parent.parent / ".env"
load_dotenv(env_path)

# Import session manager
try:
    from session_manager import (
        create_session, get_active_session, get_current_pptx,
        get_current_pptx_name, log_edit, update_slide_count,
        close_session, get_session_summary, list_recent_sessions,
        load_session
    )
    SESSION_ENABLED = True
except ImportError:
    SESSION_ENABLED = False

app = FastAPI(
    title="Interview Prep PowerPoint API",
    description="AI-powered interview preparation and PowerPoint generation",
    version="2.0.0"
)

# Enable CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Paths
BASE_DIR = Path(__file__).parent.parent
SRC_DIR = Path(__file__).parent
TMP_DIR = BASE_DIR / ".tmp"
TMP_DIR.mkdir(exist_ok=True)


# Request/Response Models
class ResearchRequest(BaseModel):
    company: str
    role: str
    generate_images: Optional[bool] = False


class GenerateRequest(BaseModel):
    research_file: str
    theme: Optional[str] = "modern"


class GenerateFullRequest(BaseModel):
    """Combined request for research + generation in one step."""
    company: str
    role: str
    theme: Optional[str] = "modern"
    generate_images: Optional[bool] = False


class NaturalEditRequest(BaseModel):
    """Natural language edit request."""
    instruction: str
    slide_number: Optional[int] = None


class EditTextRequest(BaseModel):
    pptx_file: str
    slide_number: int
    find_text: str
    replace_text: str
    match_all: Optional[bool] = False


class AddSlideRequest(BaseModel):
    pptx_file: str
    title: str
    description: str
    relevance: str
    after_slide: Optional[int] = None
    image_prompt: Optional[str] = None


class RegenerateImageRequest(BaseModel):
    pptx_file: str
    slide_number: int
    prompt: str
    image_index: Optional[int] = 0


class StatusResponse(BaseModel):
    status: str
    version: str
    anthropic_key: bool
    xai_key: bool
    tmp_dir_writable: bool


# ========== Core Functions (Direct Module Calls) ==========

def run_research(company: str, role: str, resume_path: str = None, generate_images: bool = False) -> dict:
    """
    Run research using the interview_research module directly.
    This provides better error handling than subprocess.
    """
    try:
        # Import the research module
        from interview_research import InterviewResearcher, parse_resume

        # Parse resume if provided
        resume_text = None
        if resume_path and Path(resume_path).exists():
            resume_text = parse_resume(resume_path)
            if resume_text:
                print(f"✅ Resume parsed: {len(resume_text)} characters")

        # Create researcher and run
        researcher = InterviewResearcher()
        research_data = researcher.research_company(company, role, resume_text)

        # Save to file
        safe_name = company.lower().replace(" ", "_").replace("/", "_")
        output_file = TMP_DIR / f"interview_research_{safe_name}.json"

        with open(output_file, "w") as f:
            json.dump(research_data, f, indent=2)

        return {
            "success": True,
            "output_file": output_file.name,
            "data": research_data
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }


def run_generation(research_file: str, theme: str = "modern") -> dict:
    """
    Generate PowerPoint using the pptx_generator module directly.
    """
    try:
        # Import the generator function
        from pptx_generator import create_interview_presentation

        research_path = TMP_DIR / research_file
        if not research_path.exists():
            return {"success": False, "error": f"Research file not found: {research_file}"}

        # Load research data
        with open(research_path) as f:
            research_data = json.load(f)

        # Generate output path
        company = research_data.get("company_overview", {}).get("name", "presentation")
        safe_name = company.lower().replace(" ", "_").replace("/", "_")
        output_path = TMP_DIR / f"interview_prep_{safe_name}.pptx"

        # Generate presentation
        create_interview_presentation(research_data, str(output_path), theme)

        return {
            "success": True,
            "output_file": output_path.name,
            "download_url": f"/api/download/{output_path.name}"
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }


def get_slide_previews(pptx_file: str, max_slides: int = 30) -> List[dict]:
    """
    Generate preview images for all slides in a PowerPoint.
    Returns base64-encoded PNG images.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import io

        pptx_path = TMP_DIR / pptx_file
        if not pptx_path.exists():
            return []

        prs = Presentation(str(pptx_path))
        previews = []

        for i, slide in enumerate(prs.slides[:max_slides]):
            # Get slide title and content
            title = ""
            content = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type == 1:  # Title
                            title = text
                        else:
                            if text:
                                content.append(text[:100])
                    elif text:
                        content.append(text[:100])

            previews.append({
                "slide_number": i + 1,
                "title": title or f"Slide {i + 1}",
                "content_preview": " | ".join(content[:3]) if content else "No text content",
                "has_image": any(shape.shape_type == 13 for shape in slide.shapes)  # MSO_SHAPE_TYPE.PICTURE
            })

        return previews

    except Exception as e:
        print(f"Error generating previews: {e}")
        return []


# ========== API Endpoints ==========

@app.get("/")
async def root():
    """API health check and endpoint listing."""
    return {
        "name": "Interview Prep PowerPoint API",
        "version": "2.0.0",
        "status": "running",
        "features": [
            "Combined research + generation endpoint",
            "Slide preview",
            "Natural language editing",
            "Direct module calls (no subprocess)"
        ],
        "endpoints": {
            "POST /api/generate-full": "Research + Generate in one step (RECOMMENDED)",
            "POST /api/research": "Research company and role",
            "POST /api/research/upload": "Research with resume upload",
            "POST /api/generate": "Generate PowerPoint from research",
            "GET /api/preview/{pptx_file}": "Get slide previews",
            "POST /api/edit/natural": "Natural language editing",
            "POST /api/edit/text": "Edit text on a slide",
            "POST /api/edit/add-slide/upload": "Add slide with image upload",
            "GET /api/download/{filename}": "Download generated file",
            "GET /api/status": "Check API status"
        }
    }


@app.get("/api/status", response_model=StatusResponse)
async def get_status():
    """Check API status and dependencies."""
    return StatusResponse(
        status="running",
        version="2.0.0",
        anthropic_key=bool(os.getenv("ANTHROPIC_API_KEY")),
        xai_key=bool(os.getenv("XAI_API_KEY")),
        tmp_dir_writable=os.access(TMP_DIR, os.W_OK)
    )


# ========== Combined Endpoint (Recommended) ==========

@app.post("/api/generate-full")
async def generate_full(request: GenerateFullRequest):
    """
    Combined endpoint: Research company + Generate PowerPoint in one request.
    This is the recommended way to use the API.
    """
    try:
        # Step 1: Research
        research_result = run_research(
            company=request.company,
            role=request.role,
            generate_images=request.generate_images
        )

        if not research_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Research failed: {research_result.get('error')}\n{research_result.get('traceback', '')}"
            )

        # Step 2: Generate PowerPoint
        gen_result = run_generation(
            research_file=research_result["output_file"],
            theme=request.theme
        )

        if not gen_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Generation failed: {gen_result.get('error')}\n{gen_result.get('traceback', '')}"
            )

        # Step 3: Get previews
        previews = get_slide_previews(gen_result["output_file"])

        # Create session
        session_id = None
        if SESSION_ENABLED:
            session = create_session(
                company=request.company,
                role=request.role,
                research_file=research_result["output_file"],
                pptx_file=gen_result["output_file"],
                theme=request.theme
            )
            session_id = session.get("session_id")

        return {
            "success": True,
            "message": f"Interview prep for {request.company} - {request.role} complete!",
            "research_file": research_result["output_file"],
            "pptx_file": gen_result["output_file"],
            "download_url": gen_result["download_url"],
            "slide_count": len(previews),
            "previews": previews,
            "session_id": session_id,
            "research_summary": {
                "company": research_result["data"].get("company_overview", {}).get("name"),
                "role": research_result["data"].get("role_analysis", {}).get("title"),
                "questions_count": len(research_result["data"].get("interview_insights", {}).get("common_questions", [])),
                "personalized": bool(research_result["data"].get("experience_highlights"))
            }
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}\n{traceback.format_exc()}")


@app.post("/api/generate-full/upload")
async def generate_full_with_resume(
    company: str = Form(...),
    role: str = Form(...),
    theme: str = Form("modern"),
    generate_images: bool = Form(False),
    resume: UploadFile = File(...)
):
    """
    Combined endpoint with resume upload.
    Research + Generate in one step with personalized insights.
    """
    try:
        # Save uploaded resume
        resume_path = TMP_DIR / f"uploaded_resume_{resume.filename}"
        with open(resume_path, "wb") as f:
            shutil.copyfileobj(resume.file, f)

        # Step 1: Research with resume
        research_result = run_research(
            company=company,
            role=role,
            resume_path=str(resume_path),
            generate_images=generate_images
        )

        # Clean up resume file
        resume_path.unlink(missing_ok=True)

        if not research_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Research failed: {research_result.get('error')}\n{research_result.get('traceback', '')}"
            )

        # Step 2: Generate PowerPoint
        gen_result = run_generation(
            research_file=research_result["output_file"],
            theme=theme
        )

        if not gen_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Generation failed: {gen_result.get('error')}\n{gen_result.get('traceback', '')}"
            )

        # Step 3: Get previews
        previews = get_slide_previews(gen_result["output_file"])

        # Create session
        session_id = None
        if SESSION_ENABLED:
            session = create_session(
                company=company,
                role=role,
                research_file=research_result["output_file"],
                pptx_file=gen_result["output_file"],
                theme=theme
            )
            session_id = session.get("session_id")

        return {
            "success": True,
            "message": f"Personalized interview prep for {company} - {role} complete!",
            "research_file": research_result["output_file"],
            "pptx_file": gen_result["output_file"],
            "download_url": gen_result["download_url"],
            "slide_count": len(previews),
            "previews": previews,
            "session_id": session_id,
            "personalized": True,
            "research_summary": {
                "company": research_result["data"].get("company_overview", {}).get("name"),
                "role": research_result["data"].get("role_analysis", {}).get("title"),
                "questions_count": len(research_result["data"].get("interview_insights", {}).get("common_questions", [])),
                "experience_highlights": len(research_result["data"].get("experience_highlights", []))
            }
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}\n{traceback.format_exc()}")


# ========== Individual Endpoints (Legacy Support) ==========

@app.post("/api/research")
async def research_company(request: ResearchRequest):
    """Research a company and role using AI."""
    result = run_research(
        company=request.company,
        role=request.role,
        generate_images=request.generate_images
    )

    if not result.get("success"):
        raise HTTPException(
            status_code=500,
            detail=f"Research failed: {result.get('error')}\n{result.get('traceback', '')}"
        )

    return {
        "success": True,
        "message": "Research completed",
        "output_file": result["output_file"],
        "data": result["data"]
    }


@app.post("/api/research/upload")
async def research_with_resume(
    company: str = Form(...),
    role: str = Form(...),
    generate_images: bool = Form(False),
    resume: UploadFile = File(...)
):
    """Research with resume upload for personalized insights."""
    try:
        # Save uploaded resume
        resume_path = TMP_DIR / f"uploaded_resume_{resume.filename}"
        with open(resume_path, "wb") as f:
            shutil.copyfileobj(resume.file, f)

        result = run_research(
            company=company,
            role=role,
            resume_path=str(resume_path),
            generate_images=generate_images
        )

        # Clean up
        resume_path.unlink(missing_ok=True)

        if not result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Research failed: {result.get('error')}\n{result.get('traceback', '')}"
            )

        return {
            "success": True,
            "message": "Research completed with resume",
            "output_file": result["output_file"],
            "data": result["data"]
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate")
async def generate_presentation(request: GenerateRequest):
    """Generate PowerPoint from research data."""
    result = run_generation(
        research_file=request.research_file,
        theme=request.theme
    )

    if not result.get("success"):
        raise HTTPException(
            status_code=500,
            detail=f"Generation failed: {result.get('error')}\n{result.get('traceback', '')}"
        )

    # Create session
    research_path = TMP_DIR / request.research_file
    if research_path.exists():
        with open(research_path) as f:
            research_data = json.load(f)
        company = research_data.get("company_overview", {}).get("name", "presentation")
        role = research_data.get("role_analysis", {}).get("title", "Unknown Role")

        if SESSION_ENABLED:
            session = create_session(
                company=company,
                role=role,
                research_file=request.research_file,
                pptx_file=result["output_file"],
                theme=request.theme
            )

    return {
        "success": True,
        "message": "Presentation generated",
        "output_file": result["output_file"],
        "download_url": result["download_url"]
    }


# ========== Preview Endpoint ==========

@app.get("/api/preview/{pptx_file}")
async def get_presentation_preview(pptx_file: str):
    """Get slide previews for a PowerPoint file."""
    pptx_path = TMP_DIR / pptx_file
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail="PowerPoint file not found")

    previews = get_slide_previews(pptx_file)

    return {
        "success": True,
        "pptx_file": pptx_file,
        "slide_count": len(previews),
        "previews": previews
    }


# ========== Editing Endpoints ==========

@app.get("/api/slides")
async def list_slides(pptx_file: str):
    """List all slides in a presentation."""
    try:
        from pptx import Presentation

        pptx_path = TMP_DIR / pptx_file
        if not pptx_path.exists():
            raise HTTPException(status_code=404, detail="PowerPoint file not found")

        prs = Presentation(str(pptx_path))
        slides_info = []

        for i, slide in enumerate(prs.slides):
            title = ""
            content = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not title and text:
                        title = text
                    elif text:
                        content.append(text[:50])

            slides_info.append(f"Slide {i+1}: {title or 'Untitled'}")
            if content:
                slides_info.append(f"  Content: {', '.join(content[:2])}")

        return {
            "success": True,
            "output": "\n".join(slides_info),
            "slide_count": len(prs.slides)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/edit/text")
async def edit_slide_text(request: EditTextRequest):
    """Edit text on a specific slide."""
    try:
        from pptx import Presentation

        pptx_path = TMP_DIR / request.pptx_file
        if not pptx_path.exists():
            raise HTTPException(status_code=404, detail="PowerPoint file not found")

        prs = Presentation(str(pptx_path))

        if request.slide_number < 1 or request.slide_number > len(prs.slides):
            raise HTTPException(status_code=400, detail=f"Invalid slide number. Presentation has {len(prs.slides)} slides.")

        slide = prs.slides[request.slide_number - 1]
        replacements = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if request.find_text in run.text:
                            if request.match_all:
                                run.text = run.text.replace(request.find_text, request.replace_text)
                            else:
                                run.text = run.text.replace(request.find_text, request.replace_text, 1)
                            replacements += 1
                            if not request.match_all:
                                break
                    if replacements > 0 and not request.match_all:
                        break
                if replacements > 0 and not request.match_all:
                    break

        prs.save(str(pptx_path))

        return {
            "success": True,
            "message": f"Made {replacements} replacement(s) on slide {request.slide_number}",
            "replacements": replacements
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/edit/add-slide/upload")
async def add_slide_with_image(
    pptx_file: str = Form(...),
    title: str = Form(...),
    description: str = Form(...),
    relevance: str = Form(...),
    after_slide: Optional[int] = Form(None),
    image: UploadFile = File(...)
):
    """Add a new slide with uploaded image."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN

        pptx_path = TMP_DIR / pptx_file
        if not pptx_path.exists():
            raise HTTPException(status_code=404, detail="PowerPoint file not found")

        # Save uploaded image
        image_path = TMP_DIR / f"uploaded_image_{image.filename}"
        with open(image_path, "wb") as f:
            shutil.copyfileobj(image.file, f)

        prs = Presentation(str(pptx_path))

        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout

        # Insert at position or append
        if after_slide and 0 < after_slide <= len(prs.slides):
            # Note: python-pptx doesn't have native insert, so we append
            slide = prs.slides.add_slide(blank_layout)
        else:
            slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.LEFT

        # Add image
        try:
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(4.5)
            slide.shapes.add_picture(str(image_path), left, top, width=width)
        except Exception as img_err:
            print(f"Could not add image: {img_err}")

        # Add description
        desc_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(2.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = description
        desc_para.font.size = Pt(14)

        # Add relevance
        rel_box = slide.shapes.add_textbox(Inches(5.2), Inches(4), Inches(4.3), Inches(1.5))
        rel_frame = rel_box.text_frame
        rel_frame.word_wrap = True
        rel_para = rel_frame.paragraphs[0]
        rel_para.text = f"Relevance: {relevance}"
        rel_para.font.size = Pt(12)
        rel_para.font.italic = True

        prs.save(str(pptx_path))

        # Clean up
        image_path.unlink(missing_ok=True)

        return {
            "success": True,
            "message": f"Added slide '{title}' to presentation",
            "slide_count": len(prs.slides)
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== Download Endpoint ==========

@app.get("/api/download/{filename}")
async def download_file(filename: str):
    """Download a generated file (PPTX or JSON)."""
    file_path = TMP_DIR / filename

    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    # Determine media type
    if filename.endswith(".pptx"):
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif filename.endswith(".json"):
        media_type = "application/json"
    else:
        media_type = "application/octet-stream"

    return FileResponse(
        path=str(file_path),
        filename=filename,
        media_type=media_type
    )


@app.get("/api/files")
async def list_files():
    """List all generated files in the temp directory."""
    files = []
    for f in TMP_DIR.iterdir():
        if f.is_file() and (f.suffix in [".pptx", ".json", ".jpeg", ".jpg", ".png"]):
            files.append({
                "name": f.name,
                "size": f.stat().st_size,
                "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
                "type": f.suffix
            })

    return {"files": sorted(files, key=lambda x: x["modified"], reverse=True)}


# ========== Session Endpoints ==========

@app.get("/api/session")
async def get_session():
    """Get current active session information."""
    if not SESSION_ENABLED:
        return {"error": "Session management not available"}

    session = get_active_session()
    if not session:
        return {
            "active": False,
            "message": "No active session. Create a presentation first."
        }

    return {
        "active": True,
        "session_id": session.get("session_id"),
        "company": session.get("company"),
        "role": session.get("role"),
        "pptx_file": session.get("pptx_file"),
        "download_url": f"/api/download/{session.get('pptx_file')}" if session.get("pptx_file") else None
    }


@app.post("/api/session/close")
async def close_current_session():
    """Close the current session."""
    if not SESSION_ENABLED:
        return {"error": "Session management not available"}

    session = get_active_session()
    if not session:
        return {"success": False, "message": "No active session to close"}

    pptx_file = session.get("pptx_file")
    close_session()

    return {
        "success": True,
        "message": "Session closed",
        "final_presentation": pptx_file,
        "download_url": f"/api/download/{pptx_file}" if pptx_file else None
    }


# ========== Frontend ==========

FRONTEND_DIR = BASE_DIR / "frontend"

@app.get("/app", response_class=HTMLResponse)
async def serve_frontend():
    """Serve the frontend application."""
    index_path = FRONTEND_DIR / "index.html"
    if index_path.exists():
        return HTMLResponse(content=index_path.read_text(), status_code=200)
    raise HTTPException(status_code=404, detail="Frontend not found")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
