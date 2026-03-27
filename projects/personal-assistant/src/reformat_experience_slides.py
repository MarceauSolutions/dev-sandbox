#!/usr/bin/env python3
"""
Reformat Experience Slides to Match Consistent Design

Reformats slides 14-18 to match the layout and style of slides 19-22:
- Navy blue background (0x003366)
- No accent bar
- Title at (0.55", 0.3") - 28pt white bold
- Image on left at (0.75", 1.5") - 5" wide
- Description on right at (6.55", 1.5") - white text
- "Relevance:" section at (6.55", 4.5") - white text

Usage:
    python reformat_experience_slides.py --input .tmp/file.pptx --research .tmp/research.json
    python reformat_experience_slides.py --input .tmp/file.pptx --slides 14,15,16,17,18
"""

import argparse
import sys
import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# Professional navy theme (matching slides 19-22)
NAVY_BG = RGBColor(0x00, 0x33, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def get_shape_by_type(slide, shape_type):
    """Get shapes of a specific type."""
    return [s for s in slide.shapes if s.shape_type == shape_type]


def delete_shape(shape):
    """Delete a shape from slide."""
    sp = shape._element
    sp.getparent().remove(sp)


def set_background_color(slide, color):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def reformat_slide(slide, title_text, description_text, relevance_text, keep_image=True):
    """
    Reformat a slide to match the target layout.

    Target layout (matching slides 19-22):
    - Navy background
    - Title: (0.55", 0.3"), 28pt, white, bold
    - Description: (6.55", 1.5"), 14pt, white
    - Relevance: (6.55", 4.5"), 12pt, white, starts with "Relevance:\n"
    - Image: (0.75", 1.5"), 5" wide
    """

    # Find and preserve the image
    image_blob = None
    image_ext = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_blob = shape.image.blob
                image_ext = shape.image.ext
            except:
                pass
            break

    # Delete all shapes
    shapes_to_delete = list(slide.shapes)
    for shape in shapes_to_delete:
        try:
            delete_shape(shape)
        except:
            pass

    # Set navy background
    set_background_color(slide, NAVY_BG)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.3), Inches(12.23), Inches(0.54)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Add image if we have one
    if keep_image and image_blob:
        import io
        image_stream = io.BytesIO(image_blob)
        try:
            slide.shapes.add_picture(
                image_stream,
                Inches(0.75), Inches(1.5),
                width=Inches(5.0)
            )
        except Exception as e:
            print(f"  Warning: Could not restore image: {e}")

    # Add description (right side)
    desc_box = slide.shapes.add_textbox(
        Inches(6.55), Inches(1.5), Inches(5.9), Inches(2.5)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description_text
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Add relevance section
    rel_box = slide.shapes.add_textbox(
        Inches(6.55), Inches(4.5), Inches(5.9), Inches(2.0)
    )
    tf = rel_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # Add "Relevance:" as header with text below
    p.text = "Relevance:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Add the actual relevance text as a new paragraph
    p2 = tf.add_paragraph()
    p2.text = relevance_text
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE


def extract_experience_titles():
    """Map generic experience titles to descriptive ones."""
    return {
        1: "Aerospace Oxygen Systems (Aerox)",
        2: "AW609 VTOL Qualification Testing",
        3: "NASA 747-SP Inerting System",
        4: "Engineering Documentation Review",
        5: "GD&T and CAD Expertise"
    }


def main():
    parser = argparse.ArgumentParser(
        description="Reformat experience slides to match consistent design"
    )
    parser.add_argument("--input", required=True, help="Input PowerPoint file")
    parser.add_argument("--research", help="Research JSON file for experience data")
    parser.add_argument("--slides", default="14,15,16,17,18", help="Slides to reformat (default: 14,15,16,17,18)")
    parser.add_argument("--output", help="Output file (defaults to input)")
    parser.add_argument("--open", action="store_true", help="Open file after reformatting")

    args = parser.parse_args()

    # Parse slide numbers
    slide_numbers = [int(x.strip()) for x in args.slides.split(",")]

    # Load research data if provided
    experience_data = []
    if args.research and os.path.exists(args.research):
        with open(args.research) as f:
            data = json.load(f)
            experience_data = data.get("experience_highlights", [])

    # Load presentation
    pptx_path = Path(args.input)
    if not pptx_path.exists():
        print(f"❌ File not found: {args.input}")
        return 1

    prs = Presentation(str(pptx_path))

    # Define better titles for each experience
    experience_titles = extract_experience_titles()

    print(f"\n{'='*60}")
    print(f"Reformatting Experience Slides")
    print(f"{'='*60}\n")

    for i, slide_num in enumerate(slide_numbers):
        if slide_num > len(prs.slides):
            print(f"⚠️  Slide {slide_num} out of range, skipping")
            continue

        slide = prs.slides[slide_num - 1]

        # Get experience data for this slide
        exp_index = i
        if exp_index < len(experience_data):
            exp = experience_data[exp_index]
            description = exp.get("experience", "")
            relevance = exp.get("relevance", "")
        else:
            # Extract from existing slide
            description = ""
            relevance = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if "why it matters" in text.lower():
                        relevance = text.replace("Why it matters:", "").replace("Why it matters", "").strip()
                    elif "relevance" in text.lower():
                        relevance = text.replace("Relevance:", "").replace("Relevance", "").strip()
                    elif text and "Relevant Experience" not in text and len(text) > 20:
                        if not description:
                            description = text

        # Get title
        title = experience_titles.get(i + 1, f"Experience Highlight {i + 1}")

        print(f"📝 Slide {slide_num}: {title}")

        # Reformat the slide
        reformat_slide(slide, title, description, relevance, keep_image=True)

        print(f"   ✓ Reformatted with navy background and consistent layout")

    # Save
    output_path = args.output or str(pptx_path)
    prs.save(output_path)
    print(f"\n✅ Saved to: {output_path}")

    # Open if requested
    if args.open:
        os.system(f'open "{output_path}"')

    return 0


if __name__ == "__main__":
    sys.exit(main())
