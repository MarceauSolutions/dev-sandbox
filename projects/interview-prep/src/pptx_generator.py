#!/usr/bin/env python3
"""
PowerPoint Generator Script
Creates professional PowerPoint presentations from research data.

Usage:
    python execution/pptx_generator.py --input .tmp/interview_research_apple.json
    python execution/pptx_generator.py --input .tmp/research.json --output .tmp/presentation.pptx --template modern
"""

import argparse
import json
import sys
import os
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import requests
except ImportError:
    requests = None


class PowerPointGenerator:
    """Generates professional PowerPoint presentations."""

    # Color schemes
    THEMES = {
        "modern": {
            "primary": RGBColor(0x1a, 0x1a, 0x2e),      # Dark blue
            "secondary": RGBColor(0x16, 0x21, 0x3e),    # Darker blue
            "accent": RGBColor(0xe9, 0x4d, 0x60),       # Coral red
            "text": RGBColor(0xff, 0xff, 0xff),         # White
            "text_dark": RGBColor(0x1a, 0x1a, 0x2e),    # Dark blue
            "background": RGBColor(0xff, 0xff, 0xff),   # White
        },
        "professional": {
            "primary": RGBColor(0x00, 0x33, 0x66),      # Navy
            "secondary": RGBColor(0x00, 0x52, 0x99),    # Blue
            "accent": RGBColor(0xff, 0x99, 0x00),       # Orange
            "text": RGBColor(0xff, 0xff, 0xff),         # White
            "text_dark": RGBColor(0x33, 0x33, 0x33),    # Dark gray
            "background": RGBColor(0xf5, 0xf5, 0xf5),   # Light gray
        },
        "minimal": {
            "primary": RGBColor(0x2c, 0x3e, 0x50),      # Dark slate
            "secondary": RGBColor(0x34, 0x49, 0x5e),    # Slate
            "accent": RGBColor(0x27, 0xae, 0x60),       # Green
            "text": RGBColor(0xff, 0xff, 0xff),         # White
            "text_dark": RGBColor(0x2c, 0x3e, 0x50),    # Dark slate
            "background": RGBColor(0xff, 0xff, 0xff),   # White
        }
    }

    def __init__(self, theme: str = "modern"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.theme = self.THEMES.get(theme, self.THEMES["modern"])

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add a title slide with company name and role."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Background
        self._add_background(slide, self.theme["primary"])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.theme["text"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.theme["accent"]
            p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%B %d, %Y")
        p.font.size = Pt(14)
        p.font.color.rgb = self.theme["text"]
        p.alignment = PP_ALIGN.CENTER

    def add_section_slide(self, title: str, number: int = None):
        """Add a section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        self._add_background(slide, self.theme["secondary"])

        # Section number
        if number:
            num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"0{number}" if number < 10 else str(number)
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent"]
            p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.theme["text"]
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title: str, bullets: list, subtitle: str = None):
        """Add a content slide with bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        self._add_background(slide, self.theme["background"])

        # Accent bar at top
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme["primary"]

        # Subtitle if provided
        start_y = 1.3
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.5))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = self.theme["secondary"]
            start_y = 1.9

        # Bullet points
        bullet_box = slide.shapes.add_textbox(Inches(0.75), Inches(start_y), Inches(11.833), Inches(5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets[:8]):  # Max 8 bullets per slide
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.theme["text_dark"]
            p.space_after = Pt(12)

    def add_two_column_slide(self, title: str, left_title: str, left_items: list,
                             right_title: str, right_items: list):
        """Add a slide with two columns."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        self._add_background(slide, self.theme["background"])

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme["primary"]

        # Left column title
        left_title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.theme["accent"]

        # Left column content
        left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme["text_dark"]
            p.space_after = Pt(8)

        # Right column title
        right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.theme["accent"]

        # Right column content
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme["text_dark"]
            p.space_after = Pt(8)

    def add_qa_slide(self, title: str, qa_pairs: list):
        """Add a Q&A formatted slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        self._add_background(slide, self.theme["background"])

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme["primary"]

        # Q&A pairs
        y_pos = 1.5
        for qa in qa_pairs[:3]:  # Max 3 Q&A pairs per slide
            # Question
            q_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_pos), Inches(11.833), Inches(0.6))
            tf = q_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Q: {qa.get('question', qa) if isinstance(qa, dict) else qa}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme["primary"]

            # Answer/Strategy
            if isinstance(qa, dict) and "strategy" in qa:
                a_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_pos + 0.5), Inches(11.833), Inches(1))
                tf = a_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"→ {qa['strategy']}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme["text_dark"]
                y_pos += 1.8
            else:
                y_pos += 1.0

    def add_checklist_slide(self, title: str, items: list):
        """Add a checklist slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        self._add_background(slide, self.theme["background"])

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme["primary"]

        # Checklist items
        bullet_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items[:10]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"☐ {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.theme["text_dark"]
            p.space_after = Pt(10)

    def _download_image(self, url: str) -> str:
        """Download an image from URL and return local path. Returns None on failure."""
        if not requests:
            print(f"   ⚠️ requests module not available, skipping image download")
            return None

        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()

            # Create temp file with appropriate extension
            ext = ".jpeg" if ".jpeg" in url or ".jpg" in url else ".png"
            tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp_file.write(response.content)
            tmp_file.close()
            return tmp_file.name
        except Exception as e:
            print(f"   ⚠️ Could not download image: {e}")
            return None

    def add_experience_highlight_slide(self, title: str, experience: dict, slide_num: int = 1):
        """Add a slide showcasing a specific experience with optional image."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        self._add_background(slide, self.theme["background"])

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title with slide number
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title} ({slide_num})"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme["primary"]

        # Get experience details
        exp_text = experience.get("experience", str(experience)) if isinstance(experience, dict) else str(experience)
        relevance = experience.get("relevance", "") if isinstance(experience, dict) else ""
        image_url = experience.get("image_url", "") if isinstance(experience, dict) else ""

        # Try to add image - check for local path first, then URL
        image_added = False
        image_path = experience.get("image_path", "") if isinstance(experience, dict) else ""

        # Check if it's a local file path (either image_path or image_url as local path)
        local_path = image_path or image_url
        if local_path and os.path.exists(local_path):
            try:
                slide.shapes.add_picture(local_path, Inches(0.75), Inches(1.5), width=Inches(5))
                image_added = True
                print(f"   ✓ Added local image: {local_path}")
            except Exception as e:
                print(f"   ⚠️ Could not add local image to slide: {e}")
        elif image_url and image_url.startswith(('http://', 'https://')):
            # Download from URL
            local_image = self._download_image(image_url)
            if local_image:
                try:
                    slide.shapes.add_picture(local_image, Inches(0.75), Inches(1.5), width=Inches(5))
                    image_added = True
                    os.unlink(local_image)
                except Exception as e:
                    print(f"   ⚠️ Could not add downloaded image to slide: {e}")

        # If we have an image, use two-column layout
        if image_added:
            # Right column: Experience text
            exp_box = slide.shapes.add_textbox(Inches(6.25), Inches(1.5), Inches(6.333), Inches(2.5))
            tf = exp_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = exp_text
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_dark"]

            # Relevance
            if relevance:
                rel_box = slide.shapes.add_textbox(Inches(6.25), Inches(4.2), Inches(6.333), Inches(2))
                tf = rel_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"Why it matters: {relevance}"
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = self.theme["secondary"]
        else:
            # Full width layout without image
            exp_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(2.5))
            tf = exp_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = exp_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_dark"]

            # Relevance
            if relevance:
                rel_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.833), Inches(2))
                tf = rel_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"→ Why it matters for this role: {relevance}"
                p.font.size = Pt(18)
                p.font.italic = True
                p.font.color.rgb = self.theme["secondary"]

    def add_closing_slide(self, title: str = "Good Luck!", subtitle: str = "You've got this."):
        """Add a closing slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        self._add_background(slide, self.theme["primary"])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.theme["text"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = self.theme["accent"]
        p.alignment = PP_ALIGN.CENTER

    def _add_background(self, slide, color):
        """Add a colored background to a slide."""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = color
        background.line.fill.background()

        # Send to back
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def save(self, path: str):
        """Save the presentation to a file."""
        self.prs.save(path)


def create_interview_presentation(research_data: dict, output_path: str, theme: str = "modern"):
    """Create a complete interview preparation PowerPoint from research data."""

    gen = PowerPointGenerator(theme=theme)

    company = research_data.get("company_overview", {}).get("name", "Company")
    role = research_data.get("role_analysis", {}).get("title", "Role")

    # Slide 1: Title
    gen.add_title_slide(
        title=f"{company}",
        subtitle=f"Interview Preparation: {role}"
    )

    # Slide 2: Agenda
    gen.add_section_slide("Agenda", 1)

    # Slide 3: Company Overview
    overview = research_data.get("company_overview", {})
    gen.add_content_slide(
        title="Company Overview",
        subtitle=f"{overview.get('industry', '')} | Founded {overview.get('founded', 'N/A')} | {overview.get('employees', 'N/A')} employees",
        bullets=[
            f"Headquarters: {overview.get('headquarters', 'N/A')}",
            f"Mission: {overview.get('mission', 'N/A')}",
            *[f"Product: {p}" for p in overview.get('key_products_services', [])[:3]],
        ]
    )

    # Slide 4: Recent News
    news = overview.get("recent_news", [])
    if news:
        gen.add_content_slide(
            title="Recent News & Developments",
            subtitle="Stay informed about what's happening",
            bullets=news[:6]
        )

    # Slide 5: Company Culture
    culture = research_data.get("company_culture", {})
    gen.add_two_column_slide(
        title="Company Culture",
        left_title="Core Values",
        left_items=culture.get("core_values", ["N/A"]),
        right_title="Work Environment",
        right_items=[
            culture.get("work_environment", "N/A"),
            culture.get("leadership_style", ""),
            culture.get("employee_reviews_summary", "")
        ]
    )

    # Slide 6: Role Analysis
    role_info = research_data.get("role_analysis", {})
    gen.add_content_slide(
        title=f"Role: {role}",
        subtitle=f"Department: {role_info.get('department', 'N/A')}",
        bullets=role_info.get("responsibilities", ["N/A"])
    )

    # Slide 7: Required Skills
    gen.add_two_column_slide(
        title="Skills & Success Metrics",
        left_title="Required Skills",
        left_items=role_info.get("required_skills", ["N/A"]),
        right_title="How Success is Measured",
        right_items=role_info.get("success_metrics", ["N/A"])
    )

    # Slide 8: Interview Questions (Part 1)
    interview = research_data.get("interview_insights", {})
    questions = interview.get("common_questions", [])
    if questions:
        gen.add_qa_slide(
            title="Common Interview Questions (1/2)",
            qa_pairs=questions[:3]
        )

    # Slide 9: Interview Questions (Part 2)
    if len(questions) > 3:
        gen.add_qa_slide(
            title="Common Interview Questions (2/2)",
            qa_pairs=questions[3:6]
        )

    # Slide 10: Questions to Ask
    questions_to_ask = interview.get("questions_to_ask", [])
    gen.add_content_slide(
        title="Smart Questions to Ask",
        subtitle="Show your preparation and genuine interest",
        bullets=questions_to_ask
    )

    # Slide 11: Competitive Landscape
    competitive = research_data.get("competitive_landscape", {})
    gen.add_two_column_slide(
        title="Competitive Landscape",
        left_title="Main Competitors",
        left_items=competitive.get("main_competitors", ["N/A"]),
        right_title="Industry Trends",
        right_items=competitive.get("industry_trends", ["N/A"])
    )

    # Slide 12: Your Talking Points
    talking = research_data.get("talking_points", {})
    gen.add_content_slide(
        title="Your Talking Points",
        subtitle="Key messages to communicate",
        bullets=[
            f"Why {company}: {talking.get('why_this_company', 'N/A')}",
            f"Why this role: {talking.get('why_this_role', 'N/A')}",
            f"Your value: {talking.get('value_proposition', 'N/A')}"
        ]
    )

    # Experience Highlights (if resume was provided)
    experience_highlights = research_data.get("experience_highlights", [])
    if experience_highlights and len(experience_highlights) > 0:
        # Add a section divider for experience
        candidate_name = research_data.get("candidate_name", "Your")
        gen.add_section_slide(f"{candidate_name} Relevant Experience", 2)

        # Add individual slides for top 3-5 experiences
        for i, exp in enumerate(experience_highlights[:5]):
            gen.add_experience_highlight_slide(
                title="Relevant Experience",
                experience=exp,
                slide_num=i + 1
            )

    # Preparation Checklist
    checklist = research_data.get("preparation_checklist", [])
    gen.add_checklist_slide(
        title="Preparation Checklist",
        items=checklist
    )

    # Slide 14: Closing
    gen.add_closing_slide(
        title="You're Ready!",
        subtitle=f"Go get that {role} position at {company}"
    )

    # Save
    gen.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentation from research data"
    )
    parser.add_argument(
        "--input", "-i",
        required=True,
        help="Input JSON file with research data"
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="Output PowerPoint file path"
    )
    parser.add_argument(
        "--theme", "-t",
        choices=["modern", "professional", "minimal"],
        default="modern",
        help="Presentation theme (default: modern)"
    )

    args = parser.parse_args()

    # Load research data
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"❌ Error: Input file not found: {input_path}")
        sys.exit(1)

    with open(input_path, "r") as f:
        research_data = json.load(f)

    # Set default output path
    if args.output is None:
        company = research_data.get("company_overview", {}).get("name", "presentation")
        safe_name = company.lower().replace(" ", "_").replace("/", "_")
        args.output = str(input_path.parent / f"interview_prep_{safe_name}.pptx")

    print(f"📊 Generating PowerPoint presentation...")
    print(f"   Theme: {args.theme}")

    try:
        output_path = create_interview_presentation(research_data, args.output, args.theme)
        print(f"\n✅ Presentation created successfully!")
        print(f"📁 Saved to: {output_path}")

        # Auto-open the presentation
        if sys.platform == "darwin":  # macOS
            # Close any open PowerPoint windows first
            close_script = '''
            tell application "System Events"
                if exists (processes where name is "Microsoft PowerPoint") then
                    tell application "Microsoft PowerPoint"
                        close every window saving no
                    end tell
                end if
            end tell
            '''
            try:
                subprocess.run(["osascript", "-e", close_script], capture_output=True, timeout=5)
                print(f"   Closed existing PowerPoint windows")
            except Exception:
                pass  # Ignore if PowerPoint isn't open or AppleScript fails

            # Open the new presentation
            subprocess.run(["open", output_path], check=True)
            print(f"   Opened presentation in PowerPoint")
        else:
            print(f"\n→ Open with: open \"{output_path}\"")

        return 0

    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
