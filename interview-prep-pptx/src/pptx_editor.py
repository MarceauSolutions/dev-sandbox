#!/usr/bin/env python3
"""
Interactive PowerPoint Editor
Edit existing PowerPoint presentations - modify text, replace images, add slides.

Usage:
    # List all slides and their content
    python execution/pptx_editor.py --input presentation.pptx --action list

    # Edit text on a specific slide
    python execution/pptx_editor.py --input presentation.pptx --action edit-text --slide 3 --find "old text" --replace "new text"

    # Replace an image on a slide
    python execution/pptx_editor.py --input presentation.pptx --action replace-image --slide 5 --image-index 0 --new-image /path/to/new.jpg

    # Add a new image to a slide
    python execution/pptx_editor.py --input presentation.pptx --action add-image --slide 5 --new-image /path/to/image.jpg --position "left"

    # Get detailed info about a specific slide
    python execution/pptx_editor.py --input presentation.pptx --action inspect --slide 5

    # Regenerate an experience image with new prompt
    python execution/pptx_editor.py --input presentation.pptx --action regenerate-image --slide 14 --prompt "New image description"

    # Add a new slide with AI-generated image ($0.07)
    python execution/pptx_editor.py --input presentation.pptx --action add-slide --title "My Experience" --description "Description text" --relevance "Why it matters" --prompt "AI image prompt" --after-slide 18

    # Add a new slide with YOUR OWN image (FREE)
    python execution/pptx_editor.py --input presentation.pptx --action add-slide --title "My Experience" --description "Description text" --relevance "Why it matters" --new-image /path/to/your/image.jpg --after-slide 18 --open
"""

import argparse
import json
import sys
import os
import subprocess
import tempfile
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import requests
except ImportError:
    requests = None


class PowerPointEditor:
    """Interactive editor for PowerPoint presentations."""

    def __init__(self, pptx_path: str):
        """Load a PowerPoint file for editing."""
        self.path = Path(pptx_path)
        if not self.path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

        self.prs = Presentation(str(self.path))
        self.modified = False

    def list_slides(self) -> list:
        """List all slides with their content summary."""
        slides_info = []

        for i, slide in enumerate(self.prs.slides, 1):
            slide_info = {
                "slide_number": i,
                "shapes_count": len(slide.shapes),
                "text_boxes": [],
                "images": [],
                "title": None
            }

            for shape_idx, shape in enumerate(slide.shapes):
                # Check for title
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # First text box is usually the title
                        if slide_info["title"] is None and len(text) < 100:
                            slide_info["title"] = text
                        slide_info["text_boxes"].append({
                            "index": shape_idx,
                            "text": text[:100] + "..." if len(text) > 100 else text
                        })

                # Check for images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info["images"].append({
                        "index": shape_idx,
                        "width": shape.width.inches if shape.width else 0,
                        "height": shape.height.inches if shape.height else 0,
                        "left": shape.left.inches if shape.left else 0,
                        "top": shape.top.inches if shape.top else 0
                    })

            slides_info.append(slide_info)

        return slides_info

    def inspect_slide(self, slide_number: int) -> dict:
        """Get detailed information about a specific slide."""
        if slide_number < 1 or slide_number > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_number}. Presentation has {len(self.prs.slides)} slides.")

        slide = self.prs.slides[slide_number - 1]

        slide_info = {
            "slide_number": slide_number,
            "total_slides": len(self.prs.slides),
            "shapes": []
        }

        for idx, shape in enumerate(slide.shapes):
            shape_info = {
                "index": idx,
                "type": str(shape.shape_type),
                "left": round(shape.left.inches, 2) if shape.left else 0,
                "top": round(shape.top.inches, 2) if shape.top else 0,
                "width": round(shape.width.inches, 2) if shape.width else 0,
                "height": round(shape.height.inches, 2) if shape.height else 0
            }

            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text
                shape_info["paragraphs"] = len(shape.text_frame.paragraphs)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["is_image"] = True

            slide_info["shapes"].append(shape_info)

        return slide_info

    def edit_text(self, slide_number: int, find_text: str, replace_text: str,
                  shape_index: int = None, match_all: bool = False) -> dict:
        """
        Edit text on a specific slide.

        Args:
            slide_number: 1-indexed slide number
            find_text: Text to find
            replace_text: Text to replace with
            shape_index: Optional specific shape to edit (0-indexed)
            match_all: Replace all occurrences if True

        Returns:
            Dict with results
        """
        if slide_number < 1 or slide_number > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_number}")

        slide = self.prs.slides[slide_number - 1]
        replacements = 0

        shapes_to_check = [slide.shapes[shape_index]] if shape_index is not None else slide.shapes

        for shape in shapes_to_check:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if find_text in run.text:
                            if match_all:
                                run.text = run.text.replace(find_text, replace_text)
                            else:
                                run.text = run.text.replace(find_text, replace_text, 1)
                            replacements += 1
                            if not match_all:
                                break
                    if replacements > 0 and not match_all:
                        break
                if replacements > 0 and not match_all:
                    break

        if replacements > 0:
            self.modified = True

        return {
            "success": replacements > 0,
            "replacements": replacements,
            "slide": slide_number,
            "find": find_text,
            "replace": replace_text
        }

    def replace_image(self, slide_number: int, image_index: int, new_image_path: str) -> dict:
        """
        Replace an image on a specific slide.

        Args:
            slide_number: 1-indexed slide number
            image_index: Index of the image on the slide (from inspect output)
            new_image_path: Path to new image file

        Returns:
            Dict with results
        """
        if slide_number < 1 or slide_number > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_number}")

        new_path = Path(new_image_path)
        if not new_path.exists():
            raise FileNotFoundError(f"New image not found: {new_image_path}")

        slide = self.prs.slides[slide_number - 1]

        # Find images on the slide
        images = []
        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append((idx, shape))

        if image_index >= len(images):
            raise ValueError(f"Image index {image_index} not found. Slide has {len(images)} images.")

        old_shape_idx, old_shape = images[image_index]

        # Get position and size of old image
        left = old_shape.left
        top = old_shape.top
        width = old_shape.width
        height = old_shape.height

        # Remove old image
        sp = old_shape._element
        sp.getparent().remove(sp)

        # Add new image at same position
        slide.shapes.add_picture(str(new_path), left, top, width, height)

        self.modified = True

        return {
            "success": True,
            "slide": slide_number,
            "old_image_index": image_index,
            "new_image": str(new_path),
            "position": {"left": left.inches, "top": top.inches},
            "size": {"width": width.inches, "height": height.inches}
        }

    def add_image(self, slide_number: int, image_path: str,
                  position: str = "center", width: float = 4.0) -> dict:
        """
        Add a new image to a slide.

        Args:
            slide_number: 1-indexed slide number
            image_path: Path to image file
            position: "left", "right", "center", or custom (e.g., "2.0,3.0" for left,top in inches)
            width: Width in inches (height auto-calculated to maintain aspect ratio)

        Returns:
            Dict with results
        """
        if slide_number < 1 or slide_number > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_number}")

        img_path = Path(image_path)
        if not img_path.exists():
            raise FileNotFoundError(f"Image not found: {image_path}")

        slide = self.prs.slides[slide_number - 1]

        # Calculate position based on string
        slide_width = self.prs.slide_width.inches
        slide_height = self.prs.slide_height.inches

        if position == "left":
            left = Inches(0.75)
            top = Inches(1.5)
        elif position == "right":
            left = Inches(slide_width - width - 0.75)
            top = Inches(1.5)
        elif position == "center":
            left = Inches((slide_width - width) / 2)
            top = Inches(2.0)
        elif "," in position:
            coords = position.split(",")
            left = Inches(float(coords[0]))
            top = Inches(float(coords[1]))
        else:
            left = Inches(0.75)
            top = Inches(1.5)

        # Add image
        picture = slide.shapes.add_picture(str(img_path), left, top, width=Inches(width))

        self.modified = True

        return {
            "success": True,
            "slide": slide_number,
            "image": str(img_path),
            "position": {"left": left.inches, "top": top.inches},
            "width": width
        }

    def update_slide_title(self, slide_number: int, new_title: str) -> dict:
        """Update the title of a specific slide."""
        if slide_number < 1 or slide_number > len(self.prs.slides):
            raise ValueError(f"Invalid slide number: {slide_number}")

        slide = self.prs.slides[slide_number - 1]
        old_title = None

        # Find and update title (usually first text shape)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:  # Title is usually short
                    old_title = text
                    # Update the text
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            # Preserve formatting from first run
                            first_run = paragraph.runs[0]
                            font_size = first_run.font.size
                            font_bold = first_run.font.bold
                            font_color = first_run.font.color.rgb if first_run.font.color.rgb else None

                            # Clear and set new text
                            first_run.text = new_title

                            # Remove other runs
                            for run in paragraph.runs[1:]:
                                run.text = ""

                            self.modified = True
                            return {
                                "success": True,
                                "slide": slide_number,
                                "old_title": old_title,
                                "new_title": new_title
                            }

        return {
            "success": False,
            "slide": slide_number,
            "error": "No title found on slide"
        }

    def add_experience_slide(self, title: str, description: str, relevance: str,
                               image_path: str = None, after_slide: int = None) -> dict:
        """
        Add a new experience highlight slide.

        Args:
            title: Slide title (experience name)
            description: Main experience description
            relevance: How it's relevant to the role
            image_path: Optional path to image
            after_slide: Insert after this slide number (default: end)

        Returns:
            Dict with results including new slide number
        """
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Determine where to insert
        if after_slide is None:
            insert_idx = len(self.prs.slides)
        else:
            insert_idx = after_slide

        # Create a new slide (blank layout)
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # If we need to move it to a specific position
        if after_slide is not None and after_slide < len(self.prs.slides) - 1:
            # Move slide to correct position
            slides_xml = self.prs.slides._sldIdLst
            slide_id = slides_xml[-1]  # Just added slide is at end
            slides_xml.remove(slide_id)
            slides_xml.insert(after_slide, slide_id)

        # Theme colors (matching modern theme)
        primary = RGBColor(0x1a, 0x1a, 0x2e)
        accent = RGBColor(0xe9, 0x4d, 0x60)
        text_white = RGBColor(0xff, 0xff, 0xff)

        # Add background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = primary
        background.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = text_white
        p.alignment = PP_ALIGN.LEFT

        # Add description
        desc_left = Inches(6.5) if image_path else Inches(0.75)
        desc_width = Inches(6) if image_path else Inches(11.5)
        desc_box = slide.shapes.add_textbox(desc_left, Inches(1.5), desc_width, Inches(2.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(18)
        p.font.color.rgb = text_white

        # Add relevance
        rel_box = slide.shapes.add_textbox(desc_left, Inches(4.5), desc_width, Inches(2))
        tf = rel_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Relevance:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = relevance
        p2.font.size = Pt(16)
        p2.font.color.rgb = text_white

        # Add image if provided
        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, Inches(0.75), Inches(1.5), width=Inches(5))
            except Exception as e:
                print(f"   Warning: Could not add image: {e}")

        self.modified = True

        new_slide_num = after_slide + 1 if after_slide else len(self.prs.slides)

        return {
            "success": True,
            "new_slide_number": new_slide_num,
            "title": title,
            "has_image": image_path is not None
        }

    def regenerate_image(self, slide_number: int, prompt: str, image_index: int = 0) -> dict:
        """
        Generate a new AI image and replace existing image on slide.

        Args:
            slide_number: 1-indexed slide number
            prompt: Description for new image
            image_index: Which image on the slide to replace (default: 0)

        Returns:
            Dict with results
        """
        # Import grok generator
        sys.path.insert(0, str(Path(__file__).parent))
        try:
            from grok_image_gen import GrokImageGenerator
            generator = GrokImageGenerator()
        except Exception as e:
            return {"success": False, "error": f"Could not load image generator: {e}"}

        # Generate new image
        print(f"   🎨 Generating new image: {prompt[:50]}...")
        tmp_dir = Path(__file__).parent.parent / ".tmp"
        tmp_dir.mkdir(exist_ok=True)
        output_path = str(tmp_dir / f"regenerated_slide_{slide_number}.jpeg")

        result = generator.generate_image(prompt, count=1, output_path=output_path)

        if not result:
            return {"success": False, "error": "Image generation failed"}

        # Replace the image on the slide
        try:
            replace_result = self.replace_image(slide_number, image_index, output_path)
            replace_result["prompt"] = prompt
            return replace_result
        except Exception as e:
            return {"success": False, "error": f"Failed to replace image: {e}"}

    def save(self, output_path: str = None) -> str:
        """
        Save the presentation.

        Args:
            output_path: Optional new path. If None, overwrites original.

        Returns:
            Path to saved file
        """
        save_path = output_path or str(self.path)
        self.prs.save(save_path)
        self.modified = False
        return save_path

    def save_as(self, output_path: str) -> str:
        """Save the presentation to a new file."""
        return self.save(output_path)


def print_slides_summary(slides_info: list):
    """Pretty print slides summary."""
    print(f"\n{'='*70}")
    print(f"PRESENTATION OVERVIEW ({len(slides_info)} slides)")
    print(f"{'='*70}\n")

    for slide in slides_info:
        title = slide["title"] or "(No title)"
        images = len(slide["images"])
        text_boxes = len(slide["text_boxes"])

        print(f"Slide {slide['slide_number']:2d}: {title}")
        print(f"         └─ {text_boxes} text boxes, {images} images")

        # Show image details if present
        for img in slide["images"]:
            print(f"            └─ Image[{img['index']}]: {img['width']:.1f}\" × {img['height']:.1f}\" at ({img['left']:.1f}, {img['top']:.1f})")

    print(f"\n{'='*70}\n")


def open_presentation(file_path: str):
    """Open presentation in PowerPoint, closing any existing windows first."""
    if sys.platform == "darwin":  # macOS
        # Close any open PowerPoint windows first
        close_script = '''
        tell application "System Events"
            if exists (processes where name is "Microsoft PowerPoint") then
                tell application "Microsoft PowerPoint"
                    close every window saving no
                end tell
            end if
        end tell
        '''
        try:
            subprocess.run(["osascript", "-e", close_script], capture_output=True, timeout=5)
            print(f"   Closed existing PowerPoint windows")
        except Exception:
            pass  # Ignore if PowerPoint isn't open or AppleScript fails

        # Open the new presentation
        subprocess.run(["open", file_path], check=True)
        print(f"   Opened presentation in PowerPoint")
    else:
        print(f"\n→ Open with: open \"{file_path}\"")


def print_slide_detail(slide_info: dict):
    """Pretty print detailed slide info."""
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_info['slide_number']} DETAILS")
    print(f"{'='*70}\n")

    for shape in slide_info["shapes"]:
        print(f"Shape [{shape['index']}]:")
        print(f"  Type: {shape['type']}")
        print(f"  Position: ({shape['left']}\", {shape['top']}\")")
        print(f"  Size: {shape['width']}\" × {shape['height']}\"")

        if "text" in shape:
            text_preview = shape["text"][:80] + "..." if len(shape["text"]) > 80 else shape["text"]
            print(f"  Text: \"{text_preview}\"")

        if shape.get("is_image"):
            print(f"  [IMAGE]")

        print()

    print(f"{'='*70}\n")


def main():
    parser = argparse.ArgumentParser(
        description="Interactive PowerPoint Editor - Edit presentations via command line"
    )
    parser.add_argument(
        "--input", "-i",
        required=True,
        help="Input PowerPoint file"
    )
    parser.add_argument(
        "--action", "-a",
        required=True,
        choices=["list", "inspect", "edit-text", "replace-image", "add-image",
                 "update-title", "regenerate-image", "add-slide"],
        help="Action to perform"
    )
    parser.add_argument(
        "--slide", "-s",
        type=int,
        help="Slide number (1-indexed)"
    )
    parser.add_argument(
        "--find",
        help="Text to find (for edit-text action)"
    )
    parser.add_argument(
        "--replace",
        help="Replacement text (for edit-text action)"
    )
    parser.add_argument(
        "--image-index",
        type=int,
        default=0,
        help="Index of image on slide (for replace-image, regenerate-image)"
    )
    parser.add_argument(
        "--new-image",
        help="Path to new image file"
    )
    parser.add_argument(
        "--position",
        default="left",
        help="Image position: left, right, center, or 'x,y' in inches"
    )
    parser.add_argument(
        "--width",
        type=float,
        default=5.0,
        help="Image width in inches"
    )
    parser.add_argument(
        "--title",
        help="New title text (for update-title action)"
    )
    parser.add_argument(
        "--prompt",
        help="Image generation prompt (for regenerate-image action)"
    )
    parser.add_argument(
        "--output", "-o",
        help="Output file path (default: overwrite input)"
    )
    parser.add_argument(
        "--match-all",
        action="store_true",
        help="Replace all occurrences (for edit-text)"
    )
    parser.add_argument(
        "--description",
        help="Description text (for add-slide action)"
    )
    parser.add_argument(
        "--relevance",
        help="Relevance text (for add-slide action)"
    )
    parser.add_argument(
        "--after-slide",
        type=int,
        help="Insert new slide after this slide number"
    )
    parser.add_argument(
        "--open",
        action="store_true",
        help="Open presentation after editing (closes existing PowerPoint windows)"
    )

    args = parser.parse_args()

    try:
        editor = PowerPointEditor(args.input)

        if args.action == "list":
            slides_info = editor.list_slides()
            print_slides_summary(slides_info)
            return 0

        elif args.action == "inspect":
            if not args.slide:
                print("Error: --slide required for inspect action")
                return 1
            slide_info = editor.inspect_slide(args.slide)
            print_slide_detail(slide_info)
            return 0

        elif args.action == "edit-text":
            if not all([args.slide, args.find, args.replace]):
                print("Error: --slide, --find, and --replace required for edit-text action")
                return 1
            result = editor.edit_text(args.slide, args.find, args.replace, match_all=args.match_all)

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Replaced {result['replacements']} occurrence(s)")
                print(f"   '{args.find}' → '{args.replace}'")
                print(f"📁 Saved to: {output_path}")
            else:
                print(f"⚠️ Text '{args.find}' not found on slide {args.slide}")
            return 0 if result["success"] else 1

        elif args.action == "replace-image":
            if not all([args.slide, args.new_image]):
                print("Error: --slide and --new-image required for replace-image action")
                return 1
            result = editor.replace_image(args.slide, args.image_index, args.new_image)

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Replaced image on slide {args.slide}")
                print(f"📁 Saved to: {output_path}")
            return 0 if result["success"] else 1

        elif args.action == "add-image":
            if not all([args.slide, args.new_image]):
                print("Error: --slide and --new-image required for add-image action")
                return 1
            result = editor.add_image(args.slide, args.new_image, args.position, args.width)

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Added image to slide {args.slide}")
                print(f"📁 Saved to: {output_path}")
            return 0 if result["success"] else 1

        elif args.action == "update-title":
            if not all([args.slide, args.title]):
                print("Error: --slide and --title required for update-title action")
                return 1
            result = editor.update_slide_title(args.slide, args.title)

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Updated title on slide {args.slide}")
                print(f"   '{result['old_title']}' → '{result['new_title']}'")
                print(f"📁 Saved to: {output_path}")
            else:
                print(f"⚠️ {result.get('error', 'Unknown error')}")
            return 0 if result["success"] else 1

        elif args.action == "regenerate-image":
            if not all([args.slide, args.prompt]):
                print("Error: --slide and --prompt required for regenerate-image action")
                return 1
            result = editor.regenerate_image(args.slide, args.prompt, args.image_index)

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Generated and replaced image on slide {args.slide}")
                print(f"📁 Saved to: {output_path}")

                # Open presentation if requested
                if args.open:
                    open_presentation(output_path)
            else:
                print(f"❌ {result.get('error', 'Unknown error')}")
            return 0 if result["success"] else 1

        elif args.action == "add-slide":
            if not all([args.title, args.description, args.relevance]):
                print("Error: --title, --description, and --relevance required for add-slide action")
                return 1

            # Generate image if prompt provided
            image_path = None
            if args.prompt:
                sys.path.insert(0, str(Path(__file__).parent))
                try:
                    from grok_image_gen import GrokImageGenerator
                    generator = GrokImageGenerator()
                    tmp_dir = Path(__file__).parent.parent / ".tmp"
                    tmp_dir.mkdir(exist_ok=True)
                    image_path = str(tmp_dir / f"new_slide_image.jpeg")
                    print(f"🎨 Generating image: {args.prompt[:50]}...")
                    result = generator.generate_image(args.prompt, count=1, output_path=image_path)
                    if not result:
                        print("⚠️ Image generation failed, continuing without image")
                        image_path = None
                except Exception as e:
                    print(f"⚠️ Could not generate image: {e}")
                    image_path = None
            elif args.new_image and os.path.exists(args.new_image):
                image_path = args.new_image

            result = editor.add_experience_slide(
                title=args.title,
                description=args.description,
                relevance=args.relevance,
                image_path=image_path,
                after_slide=args.after_slide
            )

            if result["success"]:
                output_path = editor.save(args.output)
                print(f"✅ Added new slide: {result['title']}")
                print(f"   New slide number: {result['new_slide_number']}")
                print(f"📁 Saved to: {output_path}")

                # Open presentation if requested
                if args.open:
                    open_presentation(output_path)
            else:
                print(f"❌ {result.get('error', 'Unknown error')}")
            return 0 if result["success"] else 1

    except FileNotFoundError as e:
        print(f"❌ File not found: {e}")
        return 1
    except ValueError as e:
        print(f"❌ Invalid value: {e}")
        return 1
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
