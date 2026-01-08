#!/usr/bin/env python3
"""
Apply Navy Theme to All Slides

Reformats all slides in a presentation to use the consistent navy theme
matching the style of slides 19-22:
- Dark navy background (#1A1A2E)
- No accent bars
- White text
- Consistent positioning

Usage:
    python apply_navy_theme.py --input .tmp/file.pptx
    python apply_navy_theme.py --input .tmp/file.pptx --slides 1,2,3,4,5
"""

import argparse
import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# Navy theme colors (matching slides 19-22)
NAVY_BG = RGBColor(0x1A, 0x1A, 0x2E)  # #1A1A2E - dark navy from reference slides
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE_ACCENT = RGBColor(0xFF, 0x99, 0x00)


def get_background_shape(slide):
    """Find the main background shape."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.width > Inches(13) and shape.height > Inches(7):
                return shape
    return None


def remove_accent_bars(slide):
    """Remove thin accent bars from slide."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Accent bars are typically full width and very thin (0.15")
            if shape.width > Inches(13) and shape.height < Inches(0.5):
                shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    return len(shapes_to_remove)


def set_background_navy(slide):
    """Set the slide background to navy blue."""
    # Try to set via background property
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_BG
    except:
        pass

    # Also update any background shape
    bg_shape = get_background_shape(slide)
    if bg_shape:
        try:
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = NAVY_BG
        except:
            pass


def set_all_text_white(slide):
    """Set all text on the slide to white."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = WHITE
                # Also set paragraph-level font color
                try:
                    paragraph.font.color.rgb = WHITE
                except:
                    pass


def reformat_slide(slide, slide_num):
    """
    Reformat a single slide to navy theme.
    Returns dict with changes made.
    """
    changes = {
        "slide": slide_num,
        "accent_bars_removed": 0,
        "background_set": False,
        "text_updated": False
    }

    # Remove accent bars
    changes["accent_bars_removed"] = remove_accent_bars(slide)

    # Set navy background
    set_background_navy(slide)
    changes["background_set"] = True

    # Set all text to white
    set_all_text_white(slide)
    changes["text_updated"] = True

    return changes


def apply_navy_theme(pptx_path: str, slide_numbers: list = None, output_path: str = None):
    """
    Apply navy theme to specified slides (or all slides).

    Args:
        pptx_path: Path to input PPTX file
        slide_numbers: List of slide numbers to reformat (1-indexed), or None for all
        output_path: Output file path (defaults to input)
    """
    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    try:
        prs = Presentation(str(path))
        total_slides = len(prs.slides)

        # Default to all slides if not specified
        if slide_numbers is None:
            slide_numbers = list(range(1, total_slides + 1))

        results = []

        for slide_num in slide_numbers:
            if slide_num < 1 or slide_num > total_slides:
                results.append({
                    "slide": slide_num,
                    "status": "skipped",
                    "reason": "out of range"
                })
                continue

            slide = prs.slides[slide_num - 1]
            changes = reformat_slide(slide, slide_num)
            changes["status"] = "updated"
            results.append(changes)

        # Save
        output = output_path or str(path)
        prs.save(output)

        return {
            "success": True,
            "output": output,
            "total_slides": total_slides,
            "slides_processed": len([r for r in results if r.get("status") == "updated"]),
            "results": results
        }

    except Exception as e:
        import traceback
        return {"success": False, "error": str(e), "traceback": traceback.format_exc()}


def main():
    parser = argparse.ArgumentParser(
        description="Apply navy theme to all slides in a presentation"
    )
    parser.add_argument("--input", required=True, help="Input PowerPoint file")
    parser.add_argument("--slides", help="Comma-separated slide numbers (default: all)")
    parser.add_argument("--output", help="Output file (defaults to input)")
    parser.add_argument("--open", action="store_true", help="Open file after processing")
    parser.add_argument("--json", action="store_true", help="Output in JSON format")

    args = parser.parse_args()

    # Parse slide numbers if provided
    slide_numbers = None
    if args.slides:
        try:
            slide_numbers = [int(x.strip()) for x in args.slides.split(",")]
        except ValueError:
            print("Error: --slides must be comma-separated numbers")
            return 1

    print(f"\n{'='*60}")
    print(f"Applying Navy Theme")
    print(f"{'='*60}\n")

    result = apply_navy_theme(args.input, slide_numbers, args.output)

    if args.json:
        print(json.dumps(result, indent=2))
    else:
        if result.get("success"):
            print(f"✅ Processed {result['slides_processed']} of {result['total_slides']} slides")
            print(f"   Output: {result['output']}")

            # Summary
            bars_removed = sum(r.get("accent_bars_removed", 0) for r in result.get("results", []))
            if bars_removed > 0:
                print(f"   Accent bars removed: {bars_removed}")

            print(f"\n📋 Changes by slide:")
            for r in result.get("results", []):
                if r.get("status") == "updated":
                    bars = f" (removed {r['accent_bars_removed']} accent bar)" if r.get("accent_bars_removed") else ""
                    print(f"   ✓ Slide {r['slide']}: Navy background, white text{bars}")
                else:
                    print(f"   ✗ Slide {r['slide']}: {r.get('reason', 'skipped')}")
        else:
            print(f"❌ Error: {result.get('error')}")
            if result.get("traceback"):
                print(result["traceback"])
            return 1

    # Open if requested
    if args.open:
        import os
        os.system(f'open "{result.get("output", args.input)}"')

    return 0


if __name__ == "__main__":
    sys.exit(main())
