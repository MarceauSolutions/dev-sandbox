#!/usr/bin/env python3
"""
Live Interactive PowerPoint Editor

Enables real-time editing of an open PowerPoint presentation.
Works alongside an open PPTX file - make changes and see them reflected immediately.

Usage:
    # Start a live editing session
    python live_editor.py --start .tmp/interview_prep_company.pptx

    # Edit text on a slide
    python live_editor.py --edit-text --slide 3 --find "old text" --replace "new text"

    # Add image to slide
    python live_editor.py --add-image --slide 14 --image .tmp/exp_img_1.jpeg

    # Refresh the open PowerPoint (triggers reload)
    python live_editor.py --refresh

    # Show current session status
    python live_editor.py --status

The workflow:
1. Open PowerPoint file manually (or use --open flag)
2. Start live session with --start
3. Make edits via commands
4. PowerPoint auto-refreshes or use --refresh to reload

Note: PowerPoint doesn't auto-reload like a browser. After edits:
- macOS: AppleScript triggers Keynote/PowerPoint to reload
- Manual: Close and reopen, or use "Revert to Saved"
"""

import argparse
import json
import sys
import os
import subprocess
import shutil
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# Session file location
SESSION_FILE = Path(".tmp/live_session.json")


def get_session() -> dict:
    """Get current live editing session."""
    if SESSION_FILE.exists():
        with open(SESSION_FILE) as f:
            return json.load(f)
    return None


def save_session(session: dict):
    """Save live editing session."""
    SESSION_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(SESSION_FILE, "w") as f:
        json.dump(session, f, indent=2)


def start_session(pptx_path: str, open_file: bool = False) -> dict:
    """
    Start a live editing session with a PowerPoint file.
    Optionally opens the file in the default application.
    """
    path = Path(pptx_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {pptx_path}"}

    # Create backup
    backup_path = path.with_suffix(".backup.pptx")
    shutil.copy(path, backup_path)

    # Load presentation to get slide count
    prs = Presentation(str(path))
    slide_count = len(prs.slides)

    session = {
        "file": str(path.absolute()),
        "file_name": path.name,
        "backup": str(backup_path.absolute()),
        "slide_count": slide_count,
        "started": datetime.now().isoformat(),
        "edits": [],
        "edit_count": 0
    }
    save_session(session)

    # Open file if requested
    if open_file:
        subprocess.run(["open", str(path)], check=False)

    return {
        "success": True,
        "message": f"Live session started for {path.name}",
        "slide_count": slide_count,
        "backup": str(backup_path),
        "file": str(path.absolute())
    }


def refresh_powerpoint() -> dict:
    """
    Trigger PowerPoint/Keynote to reload the current file.
    Uses AppleScript on macOS to close and reopen.
    """
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    file_path = session["file"]

    # AppleScript to refresh PowerPoint
    applescript_ppt = f'''
    tell application "Microsoft PowerPoint"
        if (count of presentations) > 0 then
            set currentDoc to active presentation
            set docPath to full name of currentDoc
            close currentDoc saving no
            open docPath
        end if
    end tell
    '''

    # AppleScript to refresh Keynote (if using Keynote)
    applescript_keynote = f'''
    tell application "Keynote"
        if (count of documents) > 0 then
            set currentDoc to front document
            set docPath to file of currentDoc
            close currentDoc saving no
            open docPath
        end if
    end tell
    '''

    # Try PowerPoint first, then Keynote
    try:
        # Check if PowerPoint is running
        result = subprocess.run(
            ["osascript", "-e", 'tell application "System Events" to (name of processes) contains "Microsoft PowerPoint"'],
            capture_output=True, text=True
        )

        if "true" in result.stdout.lower():
            subprocess.run(["osascript", "-e", applescript_ppt], check=False)
            return {"success": True, "message": "PowerPoint refreshed", "app": "Microsoft PowerPoint"}

        # Check Keynote
        result = subprocess.run(
            ["osascript", "-e", 'tell application "System Events" to (name of processes) contains "Keynote"'],
            capture_output=True, text=True
        )

        if "true" in result.stdout.lower():
            subprocess.run(["osascript", "-e", applescript_keynote], check=False)
            return {"success": True, "message": "Keynote refreshed", "app": "Keynote"}

        # Neither app running, just open the file
        subprocess.run(["open", file_path], check=False)
        return {"success": True, "message": "Opened file (no presentation app was running)", "app": "default"}

    except Exception as e:
        return {"success": False, "error": str(e)}


def edit_text(slide_num: int, find_text: str, replace_text: str, refresh: bool = True) -> dict:
    """Edit text on a specific slide."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    file_path = session["file"]

    try:
        prs = Presentation(file_path)

        if slide_num < 1 or slide_num > len(prs.slides):
            return {"success": False, "error": f"Slide {slide_num} out of range (1-{len(prs.slides)})"}

        slide = prs.slides[slide_num - 1]
        replacements = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_text)
                            replacements += 1

        if replacements == 0:
            return {"success": False, "error": f"Text '{find_text}' not found on slide {slide_num}"}

        prs.save(file_path)

        # Log the edit
        session["edits"].append({
            "type": "text",
            "slide": slide_num,
            "find": find_text,
            "replace": replace_text,
            "timestamp": datetime.now().isoformat()
        })
        session["edit_count"] += 1
        save_session(session)

        result = {
            "success": True,
            "message": f"Replaced {replacements} occurrence(s) on slide {slide_num}",
            "slide": slide_num,
            "replacements": replacements
        }

        # Auto-refresh if requested
        if refresh:
            refresh_result = refresh_powerpoint()
            result["refreshed"] = refresh_result.get("success", False)

        return result

    except Exception as e:
        return {"success": False, "error": str(e)}


def add_image(slide_num: int, image_path: str, position: str = "left", width: float = 4.5, refresh: bool = True) -> dict:
    """Add an image to a specific slide."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    file_path = session["file"]
    img_path = Path(image_path)

    if not img_path.exists():
        return {"success": False, "error": f"Image not found: {image_path}"}

    try:
        prs = Presentation(file_path)

        if slide_num < 1 or slide_num > len(prs.slides):
            return {"success": False, "error": f"Slide {slide_num} out of range (1-{len(prs.slides)})"}

        slide = prs.slides[slide_num - 1]

        # Calculate position based on parameter
        if position == "left":
            left = Inches(0.75)
            top = Inches(1.5)
        elif position == "right":
            left = Inches(7.5)
            top = Inches(1.5)
        elif position == "center":
            left = Inches(4.0)
            top = Inches(2.0)
        else:
            left = Inches(0.75)
            top = Inches(1.5)

        # Add image
        slide.shapes.add_picture(str(img_path), left, top, width=Inches(width))

        prs.save(file_path)

        # Log the edit
        session["edits"].append({
            "type": "image",
            "slide": slide_num,
            "image": str(img_path),
            "position": position,
            "timestamp": datetime.now().isoformat()
        })
        session["edit_count"] += 1
        save_session(session)

        result = {
            "success": True,
            "message": f"Added image to slide {slide_num}",
            "slide": slide_num,
            "image": str(img_path)
        }

        if refresh:
            refresh_result = refresh_powerpoint()
            result["refreshed"] = refresh_result.get("success", False)

        return result

    except Exception as e:
        return {"success": False, "error": str(e)}


def list_slides() -> dict:
    """List all slides in the current session."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    try:
        prs = Presentation(session["file"])
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "number": i + 1,
                "title": "",
                "has_images": False,
                "image_count": 0
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and not slide_info["title"]:
                        slide_info["title"] = text[:60]

                if shape.shape_type == 13:  # Picture
                    slide_info["has_images"] = True
                    slide_info["image_count"] += 1

            slides.append(slide_info)

        return {
            "success": True,
            "file": session["file_name"],
            "slide_count": len(slides),
            "slides": slides
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def get_status() -> dict:
    """Get current session status."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    return {
        "success": True,
        "file": session["file_name"],
        "full_path": session["file"],
        "backup": session["backup"],
        "slide_count": session["slide_count"],
        "edit_count": session["edit_count"],
        "started": session["started"],
        "recent_edits": session["edits"][-5:] if session["edits"] else []
    }


def revert_to_backup() -> dict:
    """Revert to the backup created at session start."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session. Use --start first."}

    backup_path = Path(session["backup"])
    file_path = Path(session["file"])

    if not backup_path.exists():
        return {"success": False, "error": "Backup file not found"}

    try:
        shutil.copy(backup_path, file_path)

        # Reset edit history
        session["edits"] = []
        session["edit_count"] = 0
        save_session(session)

        # Refresh
        refresh_powerpoint()

        return {
            "success": True,
            "message": f"Reverted to backup from {session['started']}"
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def end_session(keep_backup: bool = False) -> dict:
    """End the live editing session."""
    session = get_session()
    if not session:
        return {"success": False, "error": "No active session"}

    # Delete backup unless requested to keep
    if not keep_backup:
        backup_path = Path(session["backup"])
        if backup_path.exists():
            backup_path.unlink()

    # Delete session file
    if SESSION_FILE.exists():
        SESSION_FILE.unlink()

    return {
        "success": True,
        "message": "Live editing session ended",
        "edits_made": session["edit_count"]
    }


def main():
    parser = argparse.ArgumentParser(
        description="Live Interactive PowerPoint Editor",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Start a live session and open the file
  python live_editor.py --start .tmp/interview_prep.pptx --open

  # Edit text on slide 3
  python live_editor.py --edit-text --slide 3 --find "old" --replace "new"

  # Add image to slide 14
  python live_editor.py --add-image --slide 14 --image .tmp/exp_img_1.jpeg

  # List all slides
  python live_editor.py --list

  # Refresh PowerPoint to see changes
  python live_editor.py --refresh

  # Revert all changes
  python live_editor.py --revert

  # End session
  python live_editor.py --end
        """
    )

    # Session commands
    parser.add_argument("--start", metavar="FILE", help="Start live session with a PPTX file")
    parser.add_argument("--open", action="store_true", help="Open the file when starting session")
    parser.add_argument("--status", action="store_true", help="Show current session status")
    parser.add_argument("--end", action="store_true", help="End the live editing session")
    parser.add_argument("--revert", action="store_true", help="Revert to backup")

    # Edit commands
    parser.add_argument("--edit-text", action="store_true", help="Edit text on a slide")
    parser.add_argument("--add-image", action="store_true", help="Add image to a slide")
    parser.add_argument("--list", action="store_true", help="List all slides")
    parser.add_argument("--refresh", action="store_true", help="Refresh PowerPoint/Keynote")

    # Edit parameters
    parser.add_argument("--slide", type=int, help="Slide number to edit")
    parser.add_argument("--find", help="Text to find (for --edit-text)")
    parser.add_argument("--replace", help="Replacement text (for --edit-text)")
    parser.add_argument("--image", help="Image path (for --add-image)")
    parser.add_argument("--position", choices=["left", "right", "center"], default="left", help="Image position")
    parser.add_argument("--width", type=float, default=4.5, help="Image width in inches")
    parser.add_argument("--no-refresh", action="store_true", help="Don't auto-refresh after edit")

    # Output
    parser.add_argument("--json", action="store_true", help="Output in JSON format")

    args = parser.parse_args()

    result = None

    # Handle commands
    if args.start:
        result = start_session(args.start, open_file=args.open)
    elif args.status:
        result = get_status()
    elif args.end:
        result = end_session()
    elif args.revert:
        result = revert_to_backup()
    elif args.list:
        result = list_slides()
    elif args.refresh:
        result = refresh_powerpoint()
    elif args.edit_text:
        if not args.slide or not args.find or not args.replace:
            print("Error: --edit-text requires --slide, --find, and --replace")
            return 1
        result = edit_text(args.slide, args.find, args.replace, refresh=not args.no_refresh)
    elif args.add_image:
        if not args.slide or not args.image:
            print("Error: --add-image requires --slide and --image")
            return 1
        result = add_image(args.slide, args.image, args.position, args.width, refresh=not args.no_refresh)
    else:
        parser.print_help()
        return 0

    # Output result
    if result:
        if args.json:
            print(json.dumps(result, indent=2))
        else:
            if result.get("success"):
                print(f"✅ {result.get('message', 'Success')}")

                # Additional info based on command
                if "slide_count" in result and "slides" in result:
                    print(f"\n📊 {result['slide_count']} slides:")
                    for s in result["slides"]:
                        img = f" [🖼 {s['image_count']}]" if s["has_images"] else ""
                        print(f"  {s['number']:2}. {s['title']}{img}")
                elif "slide_count" in result:
                    print(f"   📊 Slides: {result['slide_count']}")
                    if "file" in result:
                        print(f"   📁 File: {result['file']}")
                    if "backup" in result:
                        print(f"   💾 Backup: {result['backup']}")
                elif "edit_count" in result:
                    print(f"   📝 Edits made: {result['edit_count']}")
                    if result.get("recent_edits"):
                        print("   📋 Recent edits:")
                        for edit in result["recent_edits"]:
                            if edit["type"] == "text":
                                print(f"      - Slide {edit['slide']}: '{edit['find']}' → '{edit['replace']}'")
                            elif edit["type"] == "image":
                                print(f"      - Slide {edit['slide']}: Added {edit['image']}")

                if result.get("refreshed"):
                    print("   🔄 PowerPoint refreshed")
            else:
                print(f"❌ {result.get('error', 'Unknown error')}")
                return 1

    return 0


if __name__ == "__main__":
    sys.exit(main())
